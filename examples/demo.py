"""Demo deck — quick run generates demo_quick.pptx.

Run from the repo root:
    python examples/demo.py

Run full deck generation:
    python examples/demo.py --full

Run Phase 3 validation:
    python examples/demo.py --phase3
"""
from __future__ import annotations

import argparse
from collections.abc import Iterator
from contextlib import contextmanager
import json
import os
import sys
import tempfile
import warnings
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from pptx import Presentation

import pptx_components as pc
import pptx_components.slide_builder as slide_builder_module


# ── Shared sample data ─────────────────────────────────────────────────────

MONTHS = ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
REVENUE = {"Revenue ($K)": [320, 380, 410, 470, 510, 580],
           "Target ($K)":  [350, 380, 400, 450, 490, 550]}
PIE_CATS = ["APAC", "EMEA", "Americas", "Other"]
PIE_VALS = [38, 27, 30, 5]
REGION_LEGEND = list(zip(PIE_CATS, [
    (59, 130, 246),
    (16, 185, 129),
    (249, 115, 22),
    (139, 92, 246),
]))

TABLE_HEADERS = ["Product", "Q1 Revenue", "Q2 Revenue", "Growth", "Status"]
TABLE_ROWS = [
    ["Alpha Suite",  "$420K", "$510K", "+21%", "On track"],
    ["Beta Platform","$280K", "$295K", "+5%",  "At risk"],
    ["Gamma API",    "$190K", "$240K", "+26%", "On track"],
    ["Delta Service","$95K",  "$88K",  "-7%",  "Behind"],
    ["Epsilon SDK",  "$60K",  "$75K",  "+25%", "On track"],
]

KPI_ITEMS = [
    ("New Leads", "1,284", "+9%", True),
    ("Conversion", "14.2%", "+1.1pp", True),
    ("Pipeline Value", "$3.4M", "+6%", True),
    ("Win Rate", "31%", "-2pp", False),
    ("Sales Cycle", "27 days", "-3 days", True),
    ("Upsell MRR", "$48K", "+12%", True),
]

TIMELINE_EVENTS = [
    ("Q1", "Discovery", "done"),
    ("Q2", "Pilot Launch", "done"),
    ("Q3", "Enterprise Rollout", "current"),
    ("Q4", "Automation Phase", "upcoming"),
    ("Q1 '27", "Scale + Expansion", "upcoming"),
]

BUILD_VS_BUY_LEFT = [
    "Full control over roadmap",
    "Can optimize for existing stack",
    "Higher upfront engineering cost",
    "Longer time-to-market",
]

BUILD_VS_BUY_RIGHT = [
    "Fastest path to launch",
    "Vendor support included",
    "Recurring license cost",
    "Less flexibility on edge cases",
]

# Phase 3 data ──────────────────────────────────────────────────────────────

HEATMAP_MATRIX = [
    [87, 91, 78, 95],
    [72, 68, 81, 74],
    [95, 97, 93, 99],
    [61, 55, 70, 63],
]
HEATMAP_ROWS = ["Mobile", "Web", "API", "Reports"]
HEATMAP_COLS = ["Q1", "Q2", "Q3", "Q4"]

RANGE_SEGMENTS_LATENCY = [
    (150, "OK", "ok"),
    (300, "Warn", "warn"),
    (600, "Critical", "error"),
]

RANGE_SEGMENTS_ERROR = [
    (1.0, "OK", "ok"),
    (3.0, "Warn", "warn"),
    (10.0, "Critical", "error"),
]

RANGE_SEGMENTS_UPTIME = [
    (95.0, "Below SLA", "error"),
    (99.0, "Degraded", "warn"),
    (100.0, "Healthy", "ok"),
]

DEMO_CODE = """\
import pptx_components as pc

builder = pc.SlideBuilder(prs)
builder.add(pc.SectionHeader("My Slide", badge_text="v1.0"))
builder.skip(0.2)
builder.add_row(
    pc.MetricCard("Revenue", "$1.2M", "+12%", True),
    pc.MetricCard("Churn",   "3.1%",  "+0.2pp", False),
    h=1.5,
)
builder.add(
    pc.BarChart(months, data, title="Monthly Trend"),
    h=2.8,
)"""

WATERFALL_CATS = ["Base", "New Logos", "Upsell", "Expansion", "Churn", "Discounts"]
WATERFALL_VALS = [500, 120, 85, 40, -60, -35]

SPARKLINE_ITEMS = [
    ("Net Revenue", "$1.28M", [0.82, 0.88, 0.93, 1.05, 1.14, 1.28], "+18%", True),
    ("Pipeline", "$3.4M", [2.4, 2.6, 2.8, 3.0, 3.2, 3.4], "+12%", True),
    ("Win Rate", "31%", [27, 28, 29, 30, 30, 31], "+2pp", True),
    ("Ticket Backlog", "42", [63, 58, 54, 49, 45, 42], "-21", True),
    ("Activation", "68%", [55, 58, 60, 63, 66, 68], "+7pp", True),
    ("Cloud Spend", "$184K", [170, 176, 181, 189, 186, 184], "-3%", True),
]

RADAR_AXES = ["Reliability", "Usability", "Security", "Performance", "Extensibility"]
RADAR_SERIES = {
    "Current": [78, 82, 88, 75, 69],
    "Target": [85, 86, 92, 84, 80],
}

FUNNEL_STAGES = [
    ("Visitors", 12840, None),
    ("Qualified", 4820, None),
    ("Pipeline", 1615, None),
    ("Closed Won", 412, None),
]

GANTT_LANES = [
    (
        "Platform",
        [
            ("Observability", 0.00, 0.32, "done"),
            ("Data Contracts", 0.28, 0.62, "current"),
            ("Failover", 0.66, 0.92, "upcoming"),
        ],
    ),
    (
        "Product",
        [
            ("Admin UX", 0.08, 0.34, "done"),
            ("Usage Alerts", 0.36, 0.63, "current"),
            ("Self-Serve", 0.70, 0.95, "upcoming"),
        ],
    ),
    (
        "GTM",
        [
            ("Messaging", 0.00, 0.22, "done"),
            ("Pricing Test", 0.24, 0.52, "current"),
            ("Partner Pack", 0.58, 0.82, "at_risk"),
        ],
    ),
    (
        "Ops",
        [
            ("SLA Review", 0.12, 0.28, "done"),
            ("Headcount Plan", 0.31, 0.57, "current"),
            ("Vendor Audit", 0.62, 0.88, "upcoming"),
        ],
    ),
]

LONG_NARRATIVE_TEXT = """Q2 opened with concentrated enterprise exposure: two strategic accounts represented 28% of net new ARR and introduced procurement timing risk. The commercial team reduced dependency by improving mid-market coverage, tightening qualification standards, and rebalancing outbound sequences toward faster-velocity segments.

Execution improved in the operating model. Onboarding cycle time dropped 19% after workflow simplification and proactive implementation playbooks. At the same time, migration volume increased support load, highlighting where product instrumentation and self-service diagnostics were still too shallow.

Financially, gross margin improved by 2.3 points due to infrastructure right-sizing and contract renegotiation in third-party services. This gain is expected to partially normalize in Q3 while regional failover hardening and observability upgrades are in flight.

The second-half plan focuses on durable retention drivers, expansion motion enablement, and service-quality stabilization. Decision velocity and cross-functional accountability are now the central execution constraints, not strategy clarity."""


# ── Slide factory functions (theme-agnostic) ───────────────────────────────

def slide_1_title(prs: Presentation) -> None:
    """Slide 1 — Title + Section Header."""
    b = pc.SlideBuilder(prs)
    b.add(pc.TitleBlock("Q2 Business Review",
                        "Performance metrics, trends, and outlook — June 2026"))
    b.skip(0.1)
    b.add(pc.SectionHeader("Key Highlights", badge_text="CONFIDENTIAL"))
    b.skip(0.2)
    b.add(pc.SectionHeader("Executive Summary"))


def slide_2_metrics(prs: Presentation) -> None:
    """Slide 2 — Metric cards + BigStat."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Performance KPIs", badge_text="Q2 2026"))
    b.skip(0.15)
    b.add_row(
        pc.MetricCard("Revenue", "$1.28M", delta="+18%", delta_positive=True),
        pc.MetricCard("Active Users", "24,370", delta="+7%", delta_positive=True),
        pc.MetricCard("Churn Rate", "3.2%", delta="+0.4pp", delta_positive=False),
        h=1.5,
    )
    b.skip(0.2)
    b.add_row(
        pc.BigStat("98.7%", "System Uptime", description="30-day rolling average"),
        pc.BigStat("4.2s", "Avg Response Time", description="p95 latency"),
        pc.BigStat("$42", "CAC", description="Customer acquisition cost"),
        h=1.8,
    )


def slide_3_table(prs: Presentation) -> None:
    """Slide 3 — Data table."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Product Revenue Breakdown"))
    b.skip(0.1)
    b.add(
        pc.DataTable(
            TABLE_HEADERS,
            TABLE_ROWS,
            weights=[3, 1.5, 1.5, 1, 1.2],
        ),
        h=pc.DataTable(TABLE_HEADERS, TABLE_ROWS).min_height,
    )


def slide_4_charts(prs: Presentation) -> None:
    """Slide 4 — Bar + Line charts side by side."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Revenue Trend Analysis"))
    b.skip(0.1)
    b.add_row(
        pc.BarChart(MONTHS, REVENUE, title="Monthly Revenue vs Target"),
        pc.LineChart(MONTHS, REVENUE, title="Revenue Trend"),
        h=2.8,
    )
    b.skip(0.15)
    b.add(pc.PieChart(PIE_CATS, PIE_VALS, title="Revenue by Region"), h=2.6)


def slide_5_lists(prs: Presentation) -> None:
    """Slide 5 — Three list styles side by side."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Roadmap & Status"))
    b.skip(0.1)
    b.add_row(
        pc.ListBlock(
            ["Redesign onboarding flow", "Launch mobile app v2",
             "Integrate SSO providers", "Expand API docs"],
            style="bullet",
            title="Backlog",
        ),
        pc.ListBlock(
            ["Define OKRs for H2", "Security audit complete",
             "Migrate to GCP Cloud Run", "Enable A/B testing framework"],
            style="number",
            title="Priorities",
        ),
        pc.ListBlock(
            ["Deploy rate limiting", "Fix dashboard load time",
             "Update privacy policy", "Archive legacy endpoints",
             "Notify affected users"],
            style="check",
            checked=[0, 2],
            title="Checklist",
        ),
        h=2.2,
    )


def slide_6_callouts(prs: Presentation) -> None:
    """Slide 6 — Callout boxes + Quote block."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Notices & Insights"))
    b.skip(0.15)
    b.add(pc.CalloutBox("Database migration scheduled for July 4 — expect 2 min downtime.", "info"))
    b.add(pc.CalloutBox("Churn has increased 0.4pp — review retention playbook this sprint.", "warning"))
    b.add(pc.CalloutBox("Uptime target of 99.5% achieved for third consecutive month.", "success"))
    b.add(pc.CalloutBox("Payment gateway timeout rate exceeded SLA threshold on June 12.", "error"))
    b.skip(0.15)
    b.add(pc.QuoteBlock(
        "The goal is not to build more features — it's to make existing ones indispensable.",
        author="Product Review, May 2026"
    ), h=1.3)


def slide_7_progress(prs: Presentation) -> None:
    """Slide 7 — Progress bars + Status badges."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Sprint Progress", badge_text="Week 24"))
    b.skip(0.2)
    b.add(pc.ProgressBar("Authentication service migration", 85))
    b.add(pc.ProgressBar("API v3 rollout", 62))
    b.add(pc.ProgressBar("Data warehouse rebuild", 40))
    b.add(pc.ProgressBar("Mobile parity features", 91, show_pct=True))
    b.add(pc.ProgressBar("Legacy cleanup", 20, max_value=100))
    b.skip(0.25)
    b.add_row(
        pc.StatusBadge("Operational", "ok"),
        pc.StatusBadge("Degraded", "warn"),
        pc.StatusBadge("Outage", "error"),
        pc.StatusBadge("Healthy", "ok"),
        h=0.3,
    )


def slide_8_composite(prs: Presentation) -> None:
    """Slide 8 — Full composition: SectionHeader + Row(MetricCards) + Chart.

    This slide proves the composition model — all primitives, no special cases.
    """
    b = pc.SlideBuilder(prs)
    b.add(pc.TitleBlock("Executive Dashboard", "Composite layout — all components composed"))
    b.skip(0.1)
    b.add_row(
        pc.MetricCard("MRR", "$128K", delta="+12%", delta_positive=True),
        pc.MetricCard("NPS", "72", delta="+4pts", delta_positive=True),
        pc.MetricCard("CAC", "$42", delta="-8%", delta_positive=True),
        pc.MetricCard("LTV", "$890", delta="+6%", delta_positive=True),
        h=1.5,
    )
    b.skip(0.1)
    b.add(
        pc.BarChart(MONTHS,
                    {"MRR ($K)": [90, 98, 105, 112, 120, 128]},
                    title="Monthly Recurring Revenue"),
        h=2.6,
    )


def slide_9_navigation(prs: Presentation) -> None:
    """Slide 9 - React-inspired navigation patterns for storytelling."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("React Pattern Adaptations", badge_text="Tabs + Steps"))
    b.skip(0.15)
    b.add(
        pc.TabsPanel(
            ["Overview", "Analytics", "Risks", "Decisions"],
            active_index=1,
            title="TabsPanel (inspired by Radix/shadcn Tabs)",
            variant="line",
            content=(
                "Q2 analytics summary: Revenue is up 18% YoY, activation improved by 6%, "
                "and churn remains above target in two enterprise segments."
            ),
        ),
        h=1.9,
    )
    b.skip(0.25)
    b.add(
        pc.StepFlow(
            ["Discovery", "Prototype", "Validation", "Launch", "Scale"],
            current=2,
            title="StepFlow (inspired by Ant Design Steps)",
            statuses=["done", "done", "current", "pending", "pending"],
        ),
        h=1.4,
    )


def slide_10_new_primitives(prs: Presentation) -> None:
    """Slide 10 - New composable primitives inspired by React ecosystems."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Phase 1 Additions", badge_text="Image + Legend + KPIGrid"))
    b.skip(0.1)

    image_path = os.path.join(os.path.dirname(__file__), "demo_dark_slides", "slide_004.png")

    b.add_row(
        pc.ImageBlock(image_path, mode="contain", border_rgb=(148, 163, 184), border_width_pt=1.0),
        pc.Legend(REGION_LEGEND, title="Region Color Legend"),
        h=2.6,
        weights=[1.6, 1.0],
    )
    b.skip(0.15)
    b.add(pc.KPIGrid(KPI_ITEMS, cols=3), h=3.0)


def slide_10b_sparkline_cards(prs: Presentation) -> None:
    """Additional slide for SparklineCard trend summaries."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Phase 1 Additions", badge_text="SparklineCard"))
    b.skip(0.12)
    b.add_row(
        pc.SparklineCard(*SPARKLINE_ITEMS[0]),
        pc.SparklineCard(*SPARKLINE_ITEMS[1]),
        pc.SparklineCard(*SPARKLINE_ITEMS[2]),
        h=1.6,
    )
    b.skip(0.18)
    b.add_row(
        pc.SparklineCard(*SPARKLINE_ITEMS[3]),
        pc.SparklineCard(*SPARKLINE_ITEMS[4]),
        pc.SparklineCard(*SPARKLINE_ITEMS[5]),
        h=1.6,
    )
    b.skip(0.2)
    b.add(
        pc.TextCard(
            "SparklineCard combines a primary KPI, change signal, and compact trend view in one surface. "
            "Use it when a standard metric card is too static for dashboard storytelling.",
            style="muted",
        ),
        h=1.0,
    )


def slide_11_timeline_comparison(prs: Presentation) -> None:
    """Slide 11 - Timeline and comparison patterns with expanded chart mode."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Phase 2 Additions", badge_text="Timeline + Comparison"))
    b.skip(0.1)
    b.add(
        pc.Timeline(
            TIMELINE_EVENTS,
            title="Program Roadmap (MUI Timeline inspired)",
        ),
        h=1.85,
    )
    b.skip(0.15)
    b.add_row(
        pc.ComparisonPanel(
            "Build In-House",
            BUILD_VS_BUY_LEFT,
            "Buy Platform",
            BUILD_VS_BUY_RIGHT,
            title="Decision Matrix (shadcn composable style)",
        ),
        pc.BarChart(
            ["Build", "Buy", "Hybrid"],
            {
                "Delivery Speed": [55, 85, 75],
                "Flexibility": [92, 60, 80],
            },
            title="Evaluation Snapshot",
            mode="bar_clustered",
        ),
        h=2.8,
        weights=[1.45, 1.0],
    )


def slide_12_heatmap_range(prs: Presentation) -> None:
    """Slide 12 - Heatmap data grid + RangeIndicator SLA gauges."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Phase 3 Additions", badge_text="Heatmap + Range"))
    b.skip(0.1)
    b.add(
        pc.Heatmap(
            HEATMAP_MATRIX,
            HEATMAP_ROWS,
            HEATMAP_COLS,
            title="Feature Adoption Score by Platform & Quarter (%)",
            colormap="sequential",
            show_values=True,
        ),
        h=pc.Heatmap(HEATMAP_MATRIX, HEATMAP_ROWS, HEATMAP_COLS).min_height,
    )
    b.skip(0.15)
    b.add_row(
        pc.RangeIndicator("API Latency (ms)", 210, RANGE_SEGMENTS_LATENCY),
        pc.RangeIndicator("Error Rate (%)", 2.1, RANGE_SEGMENTS_ERROR),
        pc.RangeIndicator("Uptime (%)", 98.5, RANGE_SEGMENTS_UPTIME, min_value=90.0),
        h=pc.RangeIndicator("x", 0, RANGE_SEGMENTS_LATENCY).min_height,
    )


def slide_13_code_annotation(prs: Presentation) -> None:
    """Slide 13 - CodeBlock + Annotation stack side by side."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Phase 3 Additions", badge_text="Code + Annotations"))
    b.skip(0.12)
    b.add_row(
        pc.CodeBlock(DEMO_CODE, language="Python", show_line_numbers=True),
        pc.Column(
            pc.Annotation("Use SlideBuilder to compose slides with rows, columns, and grids.", style="note", pointer="left"),
            pc.Spacer(0.18),
            pc.Annotation("add_row() automatically distributes width equally.", style="highlight", pointer="left"),
            pc.Spacer(0.18),
            pc.Annotation("set h= to control the rendered height of each component.", style="info", pointer="left"),
            pc.Spacer(0.18),
            pc.Annotation("Mismatched weights cause rendering issues — validate inputs.", style="warning", pointer="left"),
        ),
        h=4.7,
        weights=[1.4, 1.0],
    )


def slide_14_waterfall(prs: Presentation) -> None:
    """Slide 14 - WaterfallChart P&L breakdown."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Revenue Bridge — H1 2026", badge_text="WaterfallChart"))
    b.skip(0.1)
    b.add(
        pc.WaterfallChart(
            WATERFALL_CATS,
            WATERFALL_VALS,
            title="Revenue Contribution by Category ($K)",
            show_total=True,
            total_label="Net Revenue",
        ),
        h=4.2,
    )


def slide_15_accordion_features(prs: Presentation) -> None:
    """Slide 15 - AccordionBlock (FAQ) and FeatureGrid (product showcase)."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Additional React Patterns", badge_text="Accordion + Features"))
    b.skip(0.15)
    
    faq_items = [
        ("How do we scale headcount?", 
         "Hiring plan targets 40% growth in engineering and product through Q4 2026. "
         "Focus on mid/senior IC roles and two new manager positions."),
        ("What's our go-to-market strategy?",
         "Product-led growth with freemium tier expansion. Expected 30% conversion uplift "
         "from self-serve trial v2 launching in July."),
        ("How do we reduce churn?",
         "Retention playbook rolled out April 2026. Early wins: NPS moved from 62→68, "
         "churn rate stabilized week-over-week."),
    ]
    b.add(
        pc.AccordionBlock(
            faq_items,
            expanded_index=0,
            title="Accordion (FAQ Pattern / MUI inspired)",
        ),
        h=2.1,
    )
    b.skip(0.25)
    b.add(
        pc.FeatureGrid(
            [
                ("⚡", "Real-time Sync", "Data synchronizes across all endpoints instantly."),
                ("🔒", "Enterprise Security", "SOC 2 Type II certified infrastructure."),
                ("📊", "Rich Analytics", "Comprehensive dashboards and custom reporting."),
                ("🌍", "Global Scale", "Multi-region deployment with 99.99% SLA."),
                ("🤝", "API-First", "Full REST + GraphQL coverage for integrations."),
                ("💬", "24/7 Support", "Dedicated account team for enterprise clients."),
            ],
            columns=3,
            title="FeatureGrid (Product Showcase Pattern)",
        ),
        h=2.2,
    )


def slide_16_scatter_plot(prs: Presentation) -> None:
    """Slide 16 - ScatterPlot component for correlation analysis."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Phase 4: Advanced Charts", badge_text="ScatterPlot"))
    b.skip(0.1)
    
    scatter_points = [
        (2.5, 75, "Low cost", (100, 200, 255), 0.18),
        (4.2, 88, "Best performer", (100, 200, 255), 0.2),
        (3.1, 72, "Mid-range", (100, 200, 255), 0.17),
        (5.8, 82, "High complexity", (255, 150, 100), 0.18),
        (1.5, 65, "Simple", (150, 255, 100), 0.16),
    ]
    
    b.add(
        pc.ScatterPlot(
            points=scatter_points,
            x_label="Implementation Complexity",
            y_label="Accuracy (%)",
            title="Algorithm Performance Landscape",
            x_range=(1.0, 6.5),
            y_range=(60, 95),
            show_grid=True,
            quadrant_labels=("Best\n(Simple & Accurate)", "High Effort", "Avoid", "Feasible"),
        ),
        h=3.5,
    )


def slide_17_grouped_table(prs: Presentation) -> None:
    """Slide 17 - GroupedTable for hierarchical data."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Phase 4: Hierarchical Data", badge_text="GroupedTable"))
    b.skip(0.1)
    
    groups_data = [
        {
            "header": ("NORTH AMERICA", "$580M", "$715M", "+23%"),
            "rows": [
                ("USA", "$500M", "$620M", "+24%"),
                ("Canada", "$80M", "$95M", "+19%"),
            ],
            "footer": ("SubtotalNA", "$580M", "$715M", "+23%"),
        },
        {
            "header": ("EMEA", "$350M", "$425M", "+21%"),
            "rows": [
                ("UK", "$200M", "$245M", "+22%"),
                ("France", "$150M", "$180M", "+20%"),
            ],
            "footer": ("SubtotalEMEA", "$350M", "$425M", "+21%"),
        },
        {
            "header": ("APAC", "$220M", "$285M", "+30%"),
            "rows": [
                ("India", "$120M", "$165M", "+38%"),
                ("Australia", "$100M", "$120M", "+20%"),
            ],
            "footer": ("SubtotalAPAC", "$220M", "$285M", "+30%"),
        },
    ]
    
    b.add(
        pc.GroupedTable(
            columns=["Region", "2025 Revenue", "2026 Revenue", "YoY Growth"],
            groups=groups_data,
            title="Revenue by Region & Country",
            column_widths=[2.0, 1.5, 1.5, 1.0],
            show_dividers=True,
        ),
        h=3.8,
    )


def slide_18_chart_extensions(prs: Presentation) -> None:
    """Slide 18 - Donut and radar chart variants."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Phase 4: Chart Extensions", badge_text="Donut + Radar"))
    b.skip(0.12)
    b.add_row(
        pc.DonutChart(
            PIE_CATS,
            PIE_VALS,
            center_label="72%\nEnterprise",
            title="Revenue Mix (Donut)",
        ),
        pc.RadarChart(
            RADAR_AXES,
            RADAR_SERIES,
            title="Capability Scorecard",
            filled=True,
        ),
        h=3.0,
    )
    b.skip(0.18)
    b.add(
        pc.TextCard(
            "These chart variants round out the standard chart set: DonutChart emphasizes a headline metric in the center, "
            "while RadarChart compares capability balance across multiple axes.",
            style="muted",
        ),
        h=1.0,
    )


def slide_19_pipeline_planning(prs: Presentation) -> None:
    """Slide 19 - Funnel and Gantt planning views."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Phase 4: Planning Views", badge_text="Funnel + Gantt"))
    b.skip(0.12)
    b.add_row(
        pc.FunnelChart(FUNNEL_STAGES, title="Demand Funnel"),
        pc.GanttChart(
            GANTT_LANES,
            title="Cross-Functional Delivery Plan",
            tick_labels=["Kickoff", "Wk 3", "Wk 6", "Wk 9", "Launch"],
        ),
        h=4.2,
        weights=[0.82, 1.18],
    )


def slide_20_narrative_patterns(prs: Presentation) -> None:
    """Slide 20 - Long-form narrative patterns (two-column + auto pagination)."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Phase 5: Long-form Storytelling", badge_text="Narrative"))
    b.skip(0.12)
    b.add(
        pc.NarrativeTwoColumnPage(
            title="Executive Narrative - Structured View",
            body=(
                "Q2 performance stabilized after concentration risk mitigation and onboarding "
                "process redesign. Margin expansion creates room for selective reinvestment, "
                "but service quality and migration complexity still require focused execution "
                "through Q3."
            ),
            summary="Narrative in one slide with decision support rail.",
            sidebar_title="Decision Lens",
            sidebar_points=[
                "Protect margin gains",
                "Reduce support load",
                "Prioritize retention levers",
                "Tighten ownership by workflow",
            ],
            sidebar_note="Next review: validate churn trend by segment.",
            page=1,
            total_pages=1,
        ),
        h=4.9,
    )

    # Demonstrate helper-driven multi-slide generation for very long narratives.
    pc.build_narrative_slides(
        prs,
        title="Executive Narrative - Auto Paginated",
        text=LONG_NARRATIVE_TEXT,
        summary="Long text automatically split into multiple slides.",
        sidebar_title="Steering Notes",
        sidebar_points=[
            "Watch enterprise concentration",
            "Keep onboarding gains sticky",
            "Fund observability upgrades",
            "Track p95 latency weekly",
        ],
        sidebar_note="Auto-generated narrative section",
        max_chars_per_page=540,
        max_paragraphs_per_page=2,
    )


def slide_21_multitheme_qa(prs: Presentation) -> None:
    """Slide 21 - side-by-side theme section overrides for visual QA."""
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Theme Cascade QA", badge_text="Section Overrides"))
    b.skip(0.12)

    base_theme = b.theme
    warm_section = {
        "SURFACE": (255, 247, 234),
        "SURFACE_ALT": (246, 231, 206),
        "TEXT_PRIMARY": (74, 44, 22),
        "TEXT_SECONDARY": (111, 72, 40),
        "ACCENT": (193, 107, 41),
        "MD": 0.18,
    }
    cool_section = {
        "SURFACE": (232, 245, 255),
        "SURFACE_ALT": (212, 232, 250),
        "TEXT_PRIMARY": (18, 43, 76),
        "TEXT_SECONDARY": (47, 77, 116),
        "ACCENT": (26, 99, 194),
        "MD": 0.18,
    }

    b.add_row(
        pc.Container(
            pc.Column(
                pc.TextCard("Warm Section", "Local theme patch: editorial warm palette.", style="default"),
                pc.MetricCard("Retention", "94%", "+2pp", True),
                gap=0.12,
            ),
            border_rgb=(198, 151, 104),
            local_theme=base_theme,
            theme_patch=warm_section,
        ),
        pc.Container(
            pc.Column(
                pc.TextCard("Cool Section", "Local theme patch: corporate blue palette.", style="default"),
                pc.MetricCard("Uptime", "99.97%", "+0.04pp", True),
                gap=0.12,
            ),
            border_rgb=(115, 153, 202),
            local_theme=base_theme,
            theme_patch=cool_section,
        ),
        h=3.8,
    )

    b.skip(0.2)
    b.add(
        pc.TextCard(
            "Both cards inherit the slide theme, then apply container-scoped token patches. "
            "Use this slide to check contrast, spacing, and callout consistency across local overrides.",
            style="muted",
        ),
        h=1.05,
    )


# ── Main ───────────────────────────────────────────────────────────────────

SLIDES = [
    slide_1_title,
    slide_2_metrics,
    slide_3_table,
    slide_4_charts,
    slide_5_lists,
    slide_6_callouts,
    slide_7_progress,
    slide_8_composite,
    slide_9_navigation,
    slide_10_new_primitives,
    slide_10b_sparkline_cards,
    slide_11_timeline_comparison,
    slide_12_heatmap_range,
    slide_13_code_annotation,
    slide_14_waterfall,
    slide_15_accordion_features,
    slide_16_scatter_plot,
    slide_17_grouped_table,
    slide_18_chart_extensions,
    slide_19_pipeline_planning,
    slide_20_narrative_patterns,
    slide_21_multitheme_qa,
]


@contextmanager
def _capture_slide_builders(validate_layout: bool) -> Iterator[list[pc.SlideBuilder]]:
    builders: list[pc.SlideBuilder] = []

    if not validate_layout:
        yield builders
        return

    original_public_builder = pc.SlideBuilder
    original_module_builder = slide_builder_module.SlideBuilder

    class TrackingSlideBuilder(original_module_builder):
        def __init__(self, *args, **kwargs):
            kwargs["validate"] = kwargs.get("validate", False) or validate_layout
            super().__init__(*args, **kwargs)
            builders.append(self)

    pc.SlideBuilder = TrackingSlideBuilder
    slide_builder_module.SlideBuilder = TrackingSlideBuilder
    try:
        yield builders
    finally:
        pc.SlideBuilder = original_public_builder
        slide_builder_module.SlideBuilder = original_module_builder


def _print_and_raise_for_layout_issues(
    builders: list[pc.SlideBuilder],
    *,
    validate_layout: bool,
    strict_layout: bool,
) -> None:
    if not validate_layout:
        return

    report = pc.format_layout_validation_report(builders)
    print(report)
    if strict_layout:
        pc.raise_for_layout_issues(builders, report=report)


def build_deck(
    theme: pc.Theme,
    output_path: str,
    *,
    validate_layout: bool = False,
    strict_layout: bool = False,
) -> list[pc.SlideBuilder]:
    prs = Presentation()
    prs.slide_width = __import__('pptx.util', fromlist=['Inches']).Inches(theme.SLIDE_W)
    prs.slide_height = __import__('pptx.util', fromlist=['Inches']).Inches(theme.SLIDE_H)
    pc.set_theme(theme)
    builders: list[pc.SlideBuilder] = []

    print(f"Building: {output_path}", flush=True)
    for idx, slide_fn in enumerate(SLIDES, start=1):
        print(f"  Rendering slide {idx}/{len(SLIDES)}: {slide_fn.__name__}", flush=True)
        with _capture_slide_builders(validate_layout) as created_builders:
            slide_fn(prs)
        builders.extend(created_builders)

    _print_and_raise_for_layout_issues(
        builders,
        validate_layout=validate_layout,
        strict_layout=strict_layout,
    )

    print("  Saving file...", flush=True)
    prs.save(output_path)
    print(f"Saved: {output_path}")
    return builders


def build_quick_test_deck(
    output_path: str,
    *,
    theme: pc.Theme | None = None,
    validate_layout: bool = False,
    strict_layout: bool = False,
) -> list[pc.SlideBuilder]:
    """Deterministic smoke deck for fast local validation."""
    out_dir = os.path.dirname(__file__)
    if theme is None:
        logo_path = os.path.join(out_dir, "demo_dark_slides", "slide_004.png")
        theme = pc.BrandTheme(
            accent=(24, 119, 242),
            accent_2=(222, 72, 34),
            accent_3=(20, 153, 117),
            logo_path=logo_path,
        )
    else:
        logo_path = getattr(theme, "LOGO_PATH", None) or getattr(theme, "logo_path", None)

    prs = Presentation()
    prs.slide_width = __import__('pptx.util', fromlist=['Inches']).Inches(theme.SLIDE_W)
    prs.slide_height = __import__('pptx.util', fromlist=['Inches']).Inches(theme.SLIDE_H)
    pc.set_theme(theme)

    b = pc.SlideBuilder(prs, validate=validate_layout)
    if logo_path:
        b.set_logo(str(logo_path), x=11.2, y=0.25, w=1.6)
    b.add(pc.SectionHeader("Quick Validation", badge_text="Phase 2"))
    b.skip(0.12)
    b.add(
        pc.BarChart(
            MONTHS,
            {
                "Revenue": [320, 380, 410, 470, 510, 580],
                "Target": [350, 380, 400, 450, 490, 550],
                "Forecast": [330, 395, 430, 485, 520, 600],
            },
            title="Multi-series Palette Check",
        ),
        h=3.0,
    )

    builders = [b]
    _print_and_raise_for_layout_issues(
        builders,
        validate_layout=validate_layout,
        strict_layout=strict_layout,
    )

    prs.save(output_path)
    print(f"Saved: {output_path}")
    return builders


def build_phase3_validation(
    output_path: str,
    *,
    validate_layout: bool = False,
    strict_layout: bool = False,
) -> list[pc.SlideBuilder]:
    """Exercise Phase 3 additions: BrandTheme.from_file and BG_IMAGE background."""
    out_dir = os.path.dirname(__file__)
    image_path = os.path.join(out_dir, "demo_dark_slides", "slide_004.png")

    config = {
        "bg": "#F7FAFC",
        "surface": [255, 255, 255],
        "accent": "#0C77AA",
        "accent_2": "#C2410C",
        "accent_3": [16, 185, 129],
        "bg_image": image_path if os.path.exists(image_path) else os.path.join(out_dir, "missing_phase3_bg.png"),
        "callout": {
            "info": {"fill": "#DBEAFE", "text": "#1E3A8A"},
            "success": {"fill": [220, 252, 231], "text": [22, 101, 52]},
        },
    }

    cfg_path = ""
    try:
        with tempfile.NamedTemporaryFile(
            mode="w",
            suffix=".json",
            prefix="phase3_theme_",
            dir=out_dir,
            delete=False,
            encoding="utf-8",
        ) as tmp:
            json.dump(config, tmp)
            cfg_path = tmp.name

        theme = pc.BrandTheme.from_file(cfg_path)
        if not os.path.exists(theme.BG_IMAGE or ""):
            warnings.warn(f"Phase 3 demo BG image asset missing: {theme.BG_IMAGE}")

        prs = Presentation()
        prs.slide_width = __import__('pptx.util', fromlist=['Inches']).Inches(theme.SLIDE_W)
        prs.slide_height = __import__('pptx.util', fromlist=['Inches']).Inches(theme.SLIDE_H)
        pc.set_theme(theme)

        b = pc.SlideBuilder(prs, theme=theme, validate=validate_layout)
        b.add(pc.SectionHeader("Phase 3 Validation", badge_text="from_file + BG_IMAGE"))
        b.skip(0.15)
        b.add_row(
            pc.MetricCard("Theme Loader", "OK", "JSON", True),
            pc.MetricCard("BG Image", "Enabled", "Experimental", True),
            h=1.5,
        )
        b.skip(0.15)
        b.add(pc.CalloutBox("This slide was created with BrandTheme.from_file().", "info"))

        builders = [b]
        _print_and_raise_for_layout_issues(
            builders,
            validate_layout=validate_layout,
            strict_layout=strict_layout,
        )

        prs.save(output_path)
        print(f"Saved: {output_path}")
        return builders
    finally:
        if cfg_path and os.path.exists(cfg_path):
            os.remove(cfg_path)


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate demo decks for powerpointComponents.")
    parser.add_argument("--full", action="store_true", help="Build the dark and light demo decks.")
    parser.add_argument("--phase3", action="store_true", help="Build the Phase 3 validation deck.")
    parser.add_argument(
        "--validate-layout",
        action="store_true",
        help="Enable overflow validation and print a per-slide summary.",
    )
    parser.add_argument(
        "--strict-layout",
        action="store_true",
        help="Enable layout validation and exit non-zero when any layout issues are found.",
    )
    args = parser.parse_args()

    out_dir = os.path.dirname(__file__)
    validate_layout = args.validate_layout or args.strict_layout

    try:
        if args.full:
            build_deck(
                pc.DarkTheme(),
                os.path.join(out_dir, "demo_dark.pptx"),
                validate_layout=validate_layout,
                strict_layout=args.strict_layout,
            )
            build_deck(
                pc.LightTheme(),
                os.path.join(out_dir, "demo_light.pptx"),
                validate_layout=validate_layout,
                strict_layout=args.strict_layout,
            )
        elif args.phase3:
            build_phase3_validation(
                os.path.join(out_dir, "demo_phase3.pptx"),
                validate_layout=validate_layout,
                strict_layout=args.strict_layout,
            )
        else:
            build_quick_test_deck(
                os.path.join(out_dir, "demo_quick.pptx"),
                validate_layout=validate_layout,
                strict_layout=args.strict_layout,
            )
    except pc.LayoutValidationError:
        raise SystemExit(1)

    print("Done.")


if __name__ == "__main__":
    main()
