"""Quick test of AccordionBlock and FeatureGrid components."""
from __future__ import annotations

import json
import sys
import os
import zipfile
import xml.etree.ElementTree as ET
sys.path.insert(0, os.path.dirname(__file__))

import pytest

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

import pptx_components as pc
from pptx_components.components.chart_utils import default_theme_palette


class DenseTheme(pc.BrandTheme):
    """Spacing-dense theme used to verify cascading min-height calculations."""
    MD = 0.1
    SM = 0.1


def test_theme_patch_exposes_secondary_accents_without_mutating_base():
    """Theme patch should override ACCENT_2/ACCENT_3 while preserving base theme values."""
    base = pc.BrandTheme()
    patch = {
        "ACCENT_2": (10, 20, 30),
        "ACCENT_3": (40, 50, 60),
    }

    patched = pc.apply_theme_patch(base, patch)

    assert patched.ACCENT_2 == (10, 20, 30)
    assert patched.ACCENT_3 == (40, 50, 60)
    assert base.ACCENT_2 != (10, 20, 30)
    assert base.ACCENT_3 != (40, 50, 60)


def test_default_theme_palette_orders_primary_and_secondary_accents_first():
    """Default chart palette should place ACCENT, ACCENT_2, ACCENT_3, ACCENT_SOFT first."""
    theme = pc.BrandTheme(
        accent=(1, 2, 3),
        accent_2=(4, 5, 6),
        accent_3=(7, 8, 9),
        accent_soft=(10, 11, 12),
    )

    palette = default_theme_palette(theme)

    assert palette[:4] == [
        theme.ACCENT,
        theme.ACCENT_2,
        theme.ACCENT_3,
        theme.ACCENT_SOFT,
    ]


def test_slide_builder_set_logo_missing_file_warns_and_returns_builder():
    """Missing logo path should emit a warning and allow fluent chaining."""
    prs = Presentation()
    prs.slide_width = Inches(pc.DarkTheme().SLIDE_W)
    prs.slide_height = Inches(pc.DarkTheme().SLIDE_H)

    builder = pc.SlideBuilder(prs)
    missing = "__definitely_missing_logo_file__.png"

    with pytest.warns(UserWarning):
        result = builder.set_logo(missing, x=0.2, y=0.2, w=0.5)

    assert result is builder


def test_slide_builder_set_logo_valid_file_adds_picture_shape(tmp_path):
    """Valid logo path should add one picture shape to the current slide."""
    prs = Presentation()
    prs.slide_width = Inches(pc.DarkTheme().SLIDE_W)
    prs.slide_height = Inches(pc.DarkTheme().SLIDE_H)

    builder = pc.SlideBuilder(prs)
    before = len(builder.slide.shapes)

    # Valid 1x1 PNG bytes.
    logo_path = tmp_path / "logo.png"
    logo_path.write_bytes(
        b"\x89PNG\r\n\x1a\n"
        b"\x00\x00\x00\rIHDR"
        b"\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06\x00\x00\x00"
        b"\x1f\x15\xc4\x89"
        b"\x00\x00\x00\x0cIDATx\x9cc``\x00\x00\x00\x04\x00\x01"
        b"\x0b\xe7\x02\x9b"
        b"\x00\x00\x00\x00IEND\xaeB`\x82"
    )

    result = builder.set_logo(str(logo_path), x=0.2, y=0.2, w=0.5)
    after = len(builder.slide.shapes)

    assert result is builder
    assert after == before + 1


def test_brand_theme_from_file_json_hex_colors(tmp_path):
    """BrandTheme.from_file should parse JSON hex/list colors into RGB tuples."""
    cfg_path = tmp_path / "theme.json"
    cfg_path.write_text(
        json.dumps(
            {
                "bg": "#0C77AA",
                "accent": [10, 20, 30],
                "bg_image": "hero.png",
                "callout": {
                    "info": {
                        "fill": "#112233",
                        "text": [250, 251, 252],
                    }
                },
            }
        ),
        encoding="utf-8",
    )

    theme = pc.BrandTheme.from_file(str(cfg_path))

    assert theme.BG == (12, 119, 170)
    assert theme.ACCENT == (10, 20, 30)
    assert theme.BG_IMAGE == "hero.png"
    assert theme.CALLOUT["info"] == ((17, 34, 51), (250, 251, 252))


def test_brand_template_config_tokens_load_correctly():
    """Repository brand template config should load expected key theme tokens."""
    cfg_path = os.path.join(os.path.dirname(__file__), "examples", "brand_template_config.json")

    theme = pc.BrandTheme.from_file(cfg_path)

    assert theme.BG == (11, 28, 61)
    assert theme.ACCENT == (47, 128, 237)
    assert theme.ACCENT_2 == (31, 162, 255)
    assert theme.ACCENT_3 == (0, 82, 204)
    assert theme.CALLOUT["info"] == ((31, 79, 153), (234, 241, 255))


@pytest.mark.parametrize(
    ("config_name", "expected"),
    [
        (
            "brand_template_dark.json",
            {
                "bg": (61, 90, 94),
                "surface": (70, 95, 99),
                "accent": (95, 173, 86),
                "accent_2": (126, 200, 160),
                "accent_3": (62, 128, 53),
                "callout_key": "info",
            },
        ),
        (
            "brand_template_light.json",
            {
                "bg": (255, 255, 255),
                "surface": (255, 255, 255),
                "accent": (62, 128, 53),
                "accent_2": (61, 90, 94),
                "accent_3": (45, 102, 41),
                "callout_key": "info",
            },
        ),
    ],
)
def test_brand_template_variant_configs_load_with_sane_tokens(config_name, expected):
    """Both variant brand configs should deserialize with expected palette anchors."""
    cfg_path = os.path.join(os.path.dirname(__file__), "examples", config_name)

    theme = pc.BrandTheme.from_file(cfg_path)

    assert theme.BG == expected["bg"]
    assert theme.SURFACE == expected["surface"]
    assert theme.ACCENT == expected["accent"]
    assert theme.ACCENT_2 == expected["accent_2"]
    assert theme.ACCENT_3 == expected["accent_3"]
    assert expected["callout_key"] in theme.CALLOUT

    fill, text = theme.CALLOUT[expected["callout_key"]]
    assert len(fill) == 3
    assert len(text) == 3


def test_brand_theme_from_file_yaml_or_missing_dependency(tmp_path):
    """YAML load should work when pyyaml is installed, else raise actionable ImportError."""
    cfg_path = tmp_path / "theme.yaml"
    cfg_path.write_text(
        "bg: '#102030'\n"
        "accent: [12, 34, 56]\n"
        "callout:\n"
        "  warning:\n"
        "    fill: '#334455'\n"
        "    text: [220, 221, 222]\n",
        encoding="utf-8",
    )

    try:
        import yaml  # noqa: F401
    except ImportError:
        with pytest.raises(ImportError, match="pyyaml"):
            pc.BrandTheme.from_file(str(cfg_path))
    else:
        theme = pc.BrandTheme.from_file(str(cfg_path))
        assert theme.BG == (16, 32, 48)
        assert theme.ACCENT == (12, 34, 56)
        assert theme.CALLOUT["warning"] == ((51, 68, 85), (220, 221, 222))


def test_brand_theme_from_file_unsupported_extension(tmp_path):
    """Unsupported config extensions should raise ValueError."""
    cfg_path = tmp_path / "theme.toml"
    cfg_path.write_text("bg='#ffffff'", encoding="utf-8")

    with pytest.raises(ValueError, match="Unsupported theme config extension"):
        pc.BrandTheme.from_file(str(cfg_path))


def test_slide_builder_bg_image_missing_warns_and_falls_back_to_bg_color():
    """Missing BG_IMAGE should warn and continue rendering normally."""
    prs = Presentation()
    prs.slide_width = Inches(pc.DarkTheme().SLIDE_W)
    prs.slide_height = Inches(pc.DarkTheme().SLIDE_H)

    theme = pc.BrandTheme(bg_image="__missing_background__.png")

    with pytest.warns(UserWarning, match="Background image not found"):
        builder = pc.SlideBuilder(prs, theme=theme)

    builder.add(pc.SectionHeader("Background Fallback"))
    assert len(builder.slide.shapes) >= 1


def test_slide_builder_bg_image_valid_adds_first_shape(tmp_path):
    """Valid BG_IMAGE should be inserted as the first slide shape."""
    prs = Presentation()
    prs.slide_width = Inches(pc.DarkTheme().SLIDE_W)
    prs.slide_height = Inches(pc.DarkTheme().SLIDE_H)

    image_path = tmp_path / "bg.png"
    image_path.write_bytes(
        b"\x89PNG\r\n\x1a\n"
        b"\x00\x00\x00\rIHDR"
        b"\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06\x00\x00\x00"
        b"\x1f\x15\xc4\x89"
        b"\x00\x00\x00\x0cIDATx\x9cc``\x00\x00\x00\x04\x00\x01"
        b"\x0b\xe7\x02\x9b"
        b"\x00\x00\x00\x00IEND\xaeB`\x82"
    )

    theme = pc.BrandTheme(bg_image=str(image_path))
    builder = pc.SlideBuilder(prs, theme=theme)

    assert len(builder.slide.shapes) == 1
    assert builder.slide.shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE

    builder.add(pc.SectionHeader("Layering Check"))
    assert len(builder.slide.shapes) > 1
    assert builder.slide.shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE


def test_theme_cascade_min_height_resolution():
    """min_height_for should respect the active scoped theme."""
    container = pc.Container(pc.TextCard("Title", "Body"))
    column = pc.Column(pc.TextCard("A", "B"), pc.TextCard("C", "D"))

    default_h = container.min_height_for(pc.DarkTheme())
    dense_h = container.min_height_for(DenseTheme())
    assert dense_h < default_h

    default_col_h = column.min_height_for(pc.DarkTheme())
    dense_col_h = column.min_height_for(DenseTheme())
    assert dense_col_h < default_col_h


def test_container_local_theme_patch_scope():
    """Container local theme patch should scope style + spacing to the section subtree."""
    base = pc.DarkTheme()
    patched = pc.apply_theme_patch(base, {"MD": 0.1, "SURFACE": (1, 2, 3)})

    regular = pc.Container(pc.TextCard("Title", "Body"))
    local = pc.Container(
        pc.TextCard("Title", "Body"),
        theme_patch={"MD": 0.1, "SURFACE": (1, 2, 3)},
    )

    assert patched.MD == 0.1
    assert patched.SURFACE == (1, 2, 3)
    assert local.min_height_for(base) < regular.min_height_for(base)


def test_feature_grid_does_not_emit_negative_shape_extents(tmp_path):
    """Regression: FeatureGrid should never generate negative shape/text extents."""
    prs = Presentation()
    prs.slide_width = Inches(pc.DarkTheme().SLIDE_W)
    prs.slide_height = Inches(pc.DarkTheme().SLIDE_H)
    pc.set_theme(pc.DarkTheme())

    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Feature Grid Regression"))
    b.skip(0.15)
    b.add(
        pc.FeatureGrid(
            [
                ("*", "Feature 1", "First feature description"),
                ("#", "Feature 2", "Second feature description"),
                ("+", "Feature 3", "Third feature description"),
            ],
            columns=3,
        ),
        h=1.8,
    )

    output = tmp_path / "feature_grid_regression.pptx"
    prs.save(output)

    with zipfile.ZipFile(output, "r") as zf:
        root = ET.fromstring(zf.read("ppt/slides/slide1.xml"))

    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for ext in root.findall(".//a:ext", ns):
        cx = int(ext.attrib["cx"])
        cy = int(ext.attrib["cy"])
        assert cx >= 0
        assert cy >= 0


def test_new_components():
    """Validate new navigation components render without errors."""
    prs = Presentation()
    prs.slide_width = Inches(pc.DarkTheme().SLIDE_W)
    prs.slide_height = Inches(pc.DarkTheme().SLIDE_H)
    pc.set_theme(pc.DarkTheme())

    # Slide 1: TabsPanel & StepFlow
    print("Slide 1: TabsPanel & StepFlow...", flush=True)
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("React Navigation Patterns"))
    b.skip(0.15)
    b.add(
        pc.TabsPanel(
            ["Overview", "Analytics", "Risks"],
            active_index=1,
            content="Analytics summary for Q2.",
            variant="line",
        ),
        h=1.9,
    )
    b.skip(0.2)
    b.add(
        pc.StepFlow(
            ["Step 1", "Step 2", "Step 3"],
            current=1,
            statuses=["done", "current", "pending"],
        ),
        h=1.3,
    )

    # Slide 2: AccordionBlock & FeatureGrid
    print("Slide 2: AccordionBlock & FeatureGrid...", flush=True)
    b = pc.SlideBuilder(prs)
    b.add(pc.SectionHeader("Content Organization"))
    b.skip(0.15)
    b.add(
        pc.AccordionBlock(
            [
                ("Question 1", "Answer 1 - detailed explanation here."),
                ("Question 2", "Answer 2 - more detailed explanation."),
                ("Question 3", "Answer 3 - final answer text."),
            ],
            expanded_index=0,
        ),
        h=2.0,
    )
    b.skip(0.2)
    b.add(
        pc.FeatureGrid(
            [
                ("⚡", "Feature 1", "First feature description"),
                ("🔒", "Feature 2", "Second feature description"),
                ("📊", "Feature 3", "Third feature description"),
            ],
            columns=3,
        ),
        h=1.8,
    )

    # Save
    output = os.path.join(os.path.dirname(__file__), "test_new_components.pptx")
    print(f"Saving: {output}", flush=True)
    prs.save(output)
    print("✓ All components rendered successfully!", flush=True)


if __name__ == "__main__":
    test_new_components()
