from __future__ import annotations

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptx_components.base import Component, _resolve, add_rect, add_text_box, add_accent_bar
from pptx_components.layout import Container
from pptx_components.theme import Theme


_CALLOUT_PREFIX = {
    "info":    "INFO",
    "warning": "⚠  WARNING",
    "success": "✓  SUCCESS",
    "error":   "✕  ERROR",
}


class CalloutBox(Component):
    """Styled notification box with semantic color and a bold prefix label.

    Color pairs (fill, text) come from theme.CALLOUT — no color logic here.

    Args:
        text: The message to display.
        style: "info" | "warning" | "success" | "error"
    """

    def __init__(self, text: str, style: str = "info"):
        if style not in _CALLOUT_PREFIX:
            raise ValueError(f"style must be one of {list(_CALLOUT_PREFIX)}; got {style!r}")
        self.text = text
        self.style = style

    @property
    def min_height(self) -> float:
        return 0.75

    def render(self, slide, x: float, y: float, width: float, height: float,
               theme: Theme | None = None) -> None:
        t = _resolve(theme)
        fill_rgb, text_rgb = t.CALLOUT[self.style]
        pad = t.SM

        # Background
        add_rect(slide, x, y, width, height, fill_rgb=fill_rgb, radius=0.04)

        # Left accent bar matches the callout's semantic fill color
        bar_w = 0.06
        add_rect(slide, x, y, bar_w, height, fill_rgb=fill_rgb)

        content_x = x + bar_w + pad
        content_w = width - bar_w - pad - t.XS

        prefix = _CALLOUT_PREFIX[self.style]
        full_text = f"{prefix}  {self.text}"

        # Use a shape so we can vertically center the text
        txBox = add_rect(slide, content_x, y, content_w, height)
        from pptx_components.base import apply_no_fill, set_font, set_text_frame_margins
        apply_no_fill(txBox)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        set_text_frame_margins(tf, t.XS, t.XS, t.XS, t.XS)
        try:
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            tf.word_wrap = True
            # Vertical centering
            from pptx.oxml.ns import qn
            txBody = txBox._element.txBody
            bodyPr = txBody.find(qn('a:bodyPr'))
            if bodyPr is not None:
                bodyPr.set('anchor', 'ctr')
        except Exception:
            pass
        run = tf.paragraphs[0].add_run()
        run.text = full_text
        set_font(run, t.BODY, bold=False, color_rgb=text_rgb, font_name="Calibri")


class QuoteBlock(Component):
    """Italic pull-quote with optional attribution.

    Built on a subtle surface background — uses Container internally.
    """

    def __init__(self, text: str, author: str | None = None):
        self.text = text
        self.author = author

    @property
    def min_height(self) -> float:
        return 1.3

    def render(self, slide, x: float, y: float, width: float, height: float,
               theme: Theme | None = None) -> None:
        t = _resolve(theme)
        pad = t.MD

        # Background
        add_rect(slide, x, y, width, height, fill_rgb=t.SURFACE, radius=0.05)

        # Opening quote mark accent bar (left side)
        bar_w = 0.06
        add_rect(slide, x, y, bar_w, height, fill_rgb=t.ACCENT_SOFT)

        content_x = x + bar_w + pad
        content_w = width - bar_w - pad * 2

        if self.author:
            quote_h = height - t.MD - 0.3
            add_text_box(slide, content_x, y + pad, content_w, quote_h,
                         f'\u201c{self.text}\u201d', t.SUBHEADING, italic=True,
                         color_rgb=t.TEXT_SECONDARY, font_name="Calibri Light",
                         word_wrap=True)
            add_text_box(slide, content_x, y + height - 0.3 - t.XS, content_w, 0.3,
                         f"\u2014 {self.author}", t.CAPTION,
                         color_rgb=t.TEXT_MUTED, font_name="Calibri")
        else:
            add_text_box(slide, content_x, y + pad, content_w, height - 2 * pad,
                         f'\u201c{self.text}\u201d', t.SUBHEADING, italic=True,
                         color_rgb=t.TEXT_SECONDARY, font_name="Calibri Light",
                         word_wrap=True)
