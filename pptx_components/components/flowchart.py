from __future__ import annotations
from collections import defaultdict, deque

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from pptx_components.base import (
    Component, _resolve,
    add_rect, add_text_box,
    set_font, set_text_frame_margins,
)
from pptx_components.theme import Theme

_SHAPE_RECT = 1      # process
_SHAPE_DIAMOND = 4   # decision
_SHAPE_OVAL = 9      # terminal
_SHAPE_ROUNDED = 5   # data
_TRIANGLE = 7        # arrowhead

_SHAPE_FOR_TYPE = {
    "process":  _SHAPE_RECT,
    "decision": _SHAPE_DIAMOND,
    "terminal": _SHAPE_OVAL,
    "data":     _SHAPE_ROUNDED,
}


class FlowchartDiagram(Component):
    """Node-edge flowchart with automatic top-down layout.

    Args:
        nodes: List of node dicts: {"id": str, "label": str,
               "type": "process"|"decision"|"terminal"|"data"}
        edges: List of edge dicts: {"from": str, "to": str, "label": str (optional)}
        node_width: Width of each node in inches.
        node_height: Height of each node in inches.
        v_gap: Vertical gap between node rows in inches.
        h_gap: Horizontal gap between nodes in the same row in inches.
    """

    def __init__(
        self,
        nodes: list[dict],
        edges: list[dict],
        node_width: float = 1.8,
        node_height: float = 0.5,
        v_gap: float = 0.4,
        h_gap: float = 0.35,
    ):
        self.nodes = nodes
        self.edges = edges
        self.node_width = node_width
        self.node_height = node_height
        self.v_gap = v_gap
        self.h_gap = h_gap
        self._level, self._by_level = self._compute_levels()

    def _compute_levels(self) -> tuple[dict[str, int], dict[int, list[str]]]:
        if not self.nodes:
            return {}, {}

        all_ids = [n["id"] for n in self.nodes]
        children: dict[str, list[str]] = defaultdict(list)
        parents: dict[str, list[str]] = defaultdict(list)

        seen: set[tuple[str, str]] = set()
        for edge in self.edges:
            key = (edge["from"], edge["to"])
            if key not in seen:
                seen.add(key)
                children[edge["from"]].append(edge["to"])
                parents[edge["to"]].append(edge["from"])

        roots = [nid for nid in all_ids if not parents[nid]]
        if not roots:
            roots = [all_ids[0]]  # cycle fallback

        level: dict[str, int] = {}
        queue: deque[str] = deque()
        for r in roots:
            level[r] = 0
            queue.append(r)

        visited = set(roots)
        while queue:
            nid = queue.popleft()
            for child in children[nid]:
                new_lev = level[nid] + 1
                if child not in level or level[child] < new_lev:
                    level[child] = new_lev
                if child not in visited:
                    visited.add(child)
                    queue.append(child)

        for nid in all_ids:
            if nid not in level:
                level[nid] = 0

        by_level: dict[int, list[str]] = defaultdict(list)
        for nid in all_ids:
            by_level[level[nid]].append(nid)

        return level, dict(by_level)

    @property
    def min_height(self) -> float:
        if not self.nodes:
            return 0.5
        max_lev = max(self._level.values()) if self._level else 0
        return (max_lev + 1) * (self.node_height + self.v_gap) - self.v_gap + 0.2

    def render(
        self,
        slide,
        x: float,
        y: float,
        width: float,
        height: float,
        theme: Theme | None = None,
    ) -> None:
        if not self.nodes:
            return

        t = _resolve(theme)
        nw = self.node_width
        nh = self.node_height

        # Compute node center positions
        node_cx: dict[str, float] = {}
        node_cy: dict[str, float] = {}

        for lev, node_ids in sorted(self._by_level.items()):
            n = len(node_ids)
            for col_idx, nid in enumerate(node_ids):
                if n == 1:
                    cx = x + width / 2
                else:
                    total_w = n * nw + (n - 1) * self.h_gap
                    start_cx = x + (width - total_w) / 2 + nw / 2
                    cx = start_cx + col_idx * (nw + self.h_gap)
                cy = y + lev * (nh + self.v_gap) + nh / 2
                node_cx[nid] = cx
                node_cy[nid] = cy

        line_color = t.TEXT_MUTED
        thickness = 0.025

        # Draw connectors first (so nodes render on top)
        drawn: set[tuple[str, str]] = set()
        for edge in self.edges:
            fid, tid = edge["from"], edge["to"]
            key = (fid, tid)
            if key in drawn:
                continue
            drawn.add(key)

            if fid not in node_cx or tid not in node_cx:
                continue

            fcx, fcy = node_cx[fid], node_cy[fid]
            tcx, tcy = node_cx[tid], node_cy[tid]
            flev = self._level.get(fid, 0)
            tlev = self._level.get(tid, 0)

            # Arrowhead tip sits at the top edge of the target node
            tip_y = tcy - nh / 2

            if flev == tlev and abs(fcx - tcx) > 0.1:
                # Same level → L-shape from side of source to top of target
                sx = fcx + (nw / 2 if fcx < tcx else -nw / 2)
                sy = fcy
                ex = tcx
                ey = tip_y - thickness
                mid_y = (sy + ey) / 2
                if sy <= mid_y:
                    add_rect(slide, sx - thickness / 2, sy, thickness, mid_y - sy, fill_rgb=line_color)
                x_lo, x_hi = min(sx, ex), max(sx, ex)
                add_rect(slide, x_lo, mid_y - thickness / 2, x_hi - x_lo, thickness, fill_rgb=line_color)
                if mid_y <= ey:
                    add_rect(slide, ex - thickness / 2, mid_y, thickness, ey - mid_y, fill_rgb=line_color)
            else:
                # Different levels → vertical or L-shape from bottom of source
                sx = fcx
                sy = fcy + nh / 2
                ex = tcx
                ey = tip_y - thickness
                if abs(sx - ex) < 0.01:
                    if sy <= ey:
                        add_rect(slide, sx - thickness / 2, sy, thickness, ey - sy, fill_rgb=line_color)
                else:
                    mid_y = (sy + ey) / 2
                    if sy <= mid_y:
                        add_rect(slide, sx - thickness / 2, sy, thickness, mid_y - sy, fill_rgb=line_color)
                    x_lo, x_hi = min(sx, ex), max(sx, ex)
                    add_rect(slide, x_lo, mid_y - thickness / 2, x_hi - x_lo, thickness, fill_rgb=line_color)
                    if mid_y <= ey:
                        add_rect(slide, ex - thickness / 2, mid_y, thickness, ey - mid_y, fill_rgb=line_color)

            # Arrowhead: triangle with apex pointing down, placed at tip_y
            try:
                ps = 0.1
                tri = slide.shapes.add_shape(
                    _TRIANGLE,
                    Inches(ex - ps / 2), Inches(tip_y - ps),
                    Inches(ps), Inches(ps),
                )
                tri.fill.solid()
                tri.fill.fore_color.rgb = RGBColor(*line_color)
                tri.line.fill.background()
                tri.rotation = 180  # apex down
            except Exception:
                pass

            # Edge label at midpoint
            label_text = edge.get("label", "")
            if label_text:
                lx = (fcx + tcx) / 2 - 0.25
                ly = (fcy + tcy) / 2 - 0.14
                add_text_box(
                    slide, lx, ly, 0.5, 0.22,
                    label_text, max(8, t.CAPTION - 2),
                    color_rgb=t.TEXT_MUTED,
                    alignment=PP_ALIGN.CENTER,
                )

        # Draw nodes on top of connectors
        for n in self.nodes:
            nid = n["id"]
            label = n.get("label", nid)
            ntype = n.get("type", "process")

            if nid not in node_cx:
                continue

            cx, cy = node_cx[nid], node_cy[nid]
            nx_pos = cx - nw / 2
            ny_pos = cy - nh / 2

            shape_id = _SHAPE_FOR_TYPE.get(ntype, _SHAPE_RECT)

            try:
                shape = slide.shapes.add_shape(
                    shape_id,
                    Inches(nx_pos), Inches(ny_pos),
                    Inches(nw), Inches(nh),
                )

                if ntype == "terminal":
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(*t.ACCENT)
                    shape.line.fill.background()
                    text_color = t.BG
                elif ntype == "decision":
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(*t.ACCENT_SOFT)
                    shape.line.width = Pt(1.0)
                    shape.line.color.rgb = RGBColor(*t.ACCENT)
                    text_color = t.TEXT_PRIMARY
                else:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(*t.SURFACE)
                    shape.line.width = Pt(1.0)
                    shape.line.color.rgb = RGBColor(*t.ACCENT)
                    text_color = t.TEXT_PRIMARY

                tf = shape.text_frame
                tf.word_wrap = True
                set_text_frame_margins(tf, 0.05, 0.03, 0.05, 0.03)

                para = tf.paragraphs[0]
                para.alignment = PP_ALIGN.CENTER

                run = para.add_run()
                run.text = label
                set_font(run, t.CAPTION, bold=(ntype == "terminal"), color_rgb=text_color)

            except Exception:
                pass
