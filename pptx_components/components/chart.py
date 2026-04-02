from __future__ import annotations

import warnings

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

from pptx_components.base import Component, _resolve
from pptx_components.theme import Theme
from pptx_components.components.chart_utils import (
    chart_data_from,
    pie_data_from,
    scatter_data_from,
    default_theme_palette,
)


def _set_no_fill_xml(chart) -> None:
    """Make chart area and plot area transparent via direct XML manipulation."""
    from lxml import etree
    chart_space = chart._chartSpace

    # Chart area (c:chartSpace/c:spPr)
    spPr = chart_space.find(qn('c:spPr'))
    if spPr is None:
        spPr = etree.SubElement(chart_space, qn('c:spPr'))
    noFill = spPr.find(qn('a:noFill'))
    if noFill is None:
        # Remove any existing fill
        for child in list(spPr):
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag in ('solidFill', 'gradFill', 'pattFill', 'blipFill'):
                spPr.remove(child)
        etree.SubElement(spPr, qn('a:noFill'))

    # Plot area (c:chart/c:plotArea/c:spPr)
    c_chart = chart_space.find(qn('c:chart'))
    if c_chart is not None:
        plotArea = c_chart.find(qn('c:plotArea'))
        if plotArea is not None:
            pa_spPr = plotArea.find(qn('c:spPr'))
            if pa_spPr is None:
                pa_spPr = etree.SubElement(plotArea, qn('c:spPr'))
            noFill2 = pa_spPr.find(qn('a:noFill'))
            if noFill2 is None:
                for child in list(pa_spPr):
                    tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                    if tag in ('solidFill', 'gradFill', 'pattFill', 'blipFill'):
                        pa_spPr.remove(child)
                etree.SubElement(pa_spPr, qn('a:noFill'))


def _style_chart(chart, t: Theme, title: str | None,
                 is_line: bool = False, is_pie: bool = False) -> None:
    """Apply theme colors and optional title to a chart object."""
    # Title
    chart.has_title = title is not None
    if title:
        chart.chart_title.text_frame.text = title
        run = chart.chart_title.text_frame.paragraphs[0].runs[0]
        run.font.name = "Calibri"
        run.font.size = Pt(t.SUBHEADING)
        run.font.color.rgb = RGBColor(*t.TEXT_PRIMARY)

    # Plot area / chart area background — transparent via XML
    _set_no_fill_xml(chart)

    # Legend
    if chart.has_legend:
        chart.legend.font.color.rgb = RGBColor(*t.TEXT_SECONDARY)
        chart.legend.font.size = Pt(t.CAPTION)

    # ── Axis labels and tick marks ─────────────────────────────────────
    def _style_axis(axis, show_gridlines: bool) -> None:
        axis.tick_labels.font.color.rgb = RGBColor(*t.TEXT_SECONDARY)
        axis.tick_labels.font.size = Pt(t.CAPTION)
        axis.tick_labels.font.name = "Calibri"
        axis.format.line.color.rgb = RGBColor(*t.SURFACE_ALT)
        axis.has_major_gridlines = show_gridlines
        if show_gridlines:
            axis.major_gridlines.format.line.color.rgb = RGBColor(*t.SURFACE_ALT)

    chart_type = chart.chart_type
    is_scatter = chart_type in {
        XL_CHART_TYPE.XY_SCATTER,
        XL_CHART_TYPE.XY_SCATTER_LINES,
        XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
        XL_CHART_TYPE.XY_SCATTER_SMOOTH,
        XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS,
    }

    if is_scatter:
        try:
            # In XY charts both axes are value axes; category_axis is the X axis.
            _style_axis(chart.category_axis, show_gridlines=False)
            _style_axis(chart.value_axis, show_gridlines=True)
        except (AttributeError, ValueError):
            pass
    else:
        try:
            _style_axis(chart.category_axis, show_gridlines=False)
        except (AttributeError, ValueError):
            pass  # pie/donut/radar charts have no category axis

        try:
            _style_axis(chart.value_axis, show_gridlines=True)
        except (AttributeError, ValueError):
            pass

    # ── Series colors ──────────────────────────────────────────────────
    palette = default_theme_palette(t)

    if is_pie:
        # Color each data point with a distinct hue
        plot = chart.plots[0]
        for i, point in enumerate(plot.series[0].points):
            color = palette[i % len(palette)]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RGBColor(*color)
    else:
        for i, series in enumerate(chart.series):
            color = palette[i % len(palette)]
            if is_line:
                # Lines need visible strokes, not fills
                series.format.line.color.rgb = RGBColor(*color)
                series.format.line.width = Pt(2.5)
                if chart_type == XL_CHART_TYPE.LINE:
                    series.smooth = True
            else:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = RGBColor(*color)
                series.format.line.fill.background()


def _add_chart_shape(slide, chart_data, chart_type, x, y, w, h):
    from pptx.util import Inches
    return slide.shapes.add_chart(chart_type, Inches(x), Inches(y), Inches(w), Inches(h), chart_data)


class BarChart(Component):
    """Vertical bar chart.

    Args:
        categories: X-axis labels.
        series: Dict of series name → list of values.
        title: Optional chart title.
        stacked: Use stacked bar chart. Deprecated — use mode= instead.
        mode: Optional explicit chart mode. Supported values:
            "column_clustered" (default), "column_stacked", "column_stacked_100",
            "bar_clustered", "bar_stacked".
    """

    def __init__(self, categories: list[str], series: dict[str, list[float]],
                 title: str | None = None, stacked: bool = False,
                 mode: str | None = None):
        if mode is not None and mode not in (
            "column_clustered", "column_stacked", "column_stacked_100",
            "bar_clustered", "bar_stacked",
        ):
            raise ValueError(
                "mode must be one of: 'column_clustered', 'column_stacked', "
                "'column_stacked_100', 'bar_clustered', 'bar_stacked'"
            )
        self.categories = categories
        self.series = series
        self.title = title
        self.stacked = stacked
        self.mode = mode
        if stacked:
            warnings.warn(
                "The 'stacked' parameter is deprecated; use mode='column_stacked' or mode='bar_stacked' instead.",
                DeprecationWarning,
                stacklevel=2,
            )

    @property
    def min_height(self) -> float:
        return 2.5

    def render(self, slide, x: float, y: float, width: float, height: float,
               theme: Theme | None = None) -> None:
        t = _resolve(theme)
        cd = chart_data_from(self.categories, self.series)
        if self.mode == "column_stacked":
            chart_type = XL_CHART_TYPE.COLUMN_STACKED
        elif self.mode == "column_stacked_100":
            chart_type = XL_CHART_TYPE.COLUMN_STACKED_100
        elif self.mode == "bar_clustered":
            chart_type = XL_CHART_TYPE.BAR_CLUSTERED
        elif self.mode == "bar_stacked":
            chart_type = XL_CHART_TYPE.BAR_STACKED
        elif self.mode == "column_clustered":
            chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
        elif self.stacked:
            chart_type = XL_CHART_TYPE.COLUMN_STACKED
        else:
            chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
        graphic = _add_chart_shape(slide, cd, chart_type, x, y, width, height)
        _style_chart(graphic.chart, t, self.title, is_line=False, is_pie=False)


class LineChart(Component):
    """Line chart.

    Args:
        categories: X-axis labels.
        series: Dict of series name → list of values.
        title: Optional chart title.
    """

    def __init__(self, categories: list[str], series: dict[str, list[float]],
                 title: str | None = None):
        self.categories = categories
        self.series = series
        self.title = title

    @property
    def min_height(self) -> float:
        return 2.5

    def render(self, slide, x: float, y: float, width: float, height: float,
               theme: Theme | None = None) -> None:
        t = _resolve(theme)
        cd = chart_data_from(self.categories, self.series)
        graphic = _add_chart_shape(slide, cd, XL_CHART_TYPE.LINE, x, y, width, height)
        _style_chart(graphic.chart, t, self.title, is_line=True)


class PieChart(Component):
    """Pie chart.

    Args:
        categories: Slice labels.
        values: Slice values.
        title: Optional chart title.
    """

    def __init__(self, categories: list[str], values: list[float],
                 title: str | None = None):
        self.categories = categories
        self.values = values
        self.title = title

    @property
    def min_height(self) -> float:
        return 2.5

    def render(self, slide, x: float, y: float, width: float, height: float,
               theme: Theme | None = None) -> None:
        t = _resolve(theme)
        cd = pie_data_from(self.categories, self.values)
        graphic = _add_chart_shape(slide, cd, XL_CHART_TYPE.PIE, x, y, width, height)
        _style_chart(graphic.chart, t, self.title, is_pie=True)


class ScatterChart(Component):
    """XY scatter chart for correlation analysis and portfolio views.

    Args:
        series: Dict of series name -> list of (x, y) tuples.
        title: Optional chart title.
    """

    def __init__(
        self,
        series: dict[str, list[tuple[float, float]]],
        title: str | None = None,
    ):
        for name, points in series.items():
            if len(points) < 1:
                raise ValueError(f"Series '{name}' must have at least 1 point.")
            for pt in points:
                if (
                    not isinstance(pt, (list, tuple))
                    or len(pt) != 2
                    or not all(isinstance(v, (int, float)) for v in pt)
                ):
                    raise ValueError(
                        f"Each point in series '{name}' must be a 2-tuple of numbers;"
                        f" got {pt!r}."
                    )
        self.series = series
        self.title = title

    @property
    def min_height(self) -> float:
        return 2.5

    def render(self, slide, x: float, y: float, width: float, height: float,
               theme: Theme | None = None) -> None:
        t = _resolve(theme)
        cd = scatter_data_from(self.series)
        graphic = _add_chart_shape(slide, cd, XL_CHART_TYPE.XY_SCATTER, x, y, width, height)
        chart = graphic.chart
        _style_chart(chart, t, self.title, is_line=False, is_pie=False)

        # Apply palette colors and marker sizes per series
        palette = default_theme_palette(t)
        for i, ser in enumerate(chart.series):
            color = palette[i % len(palette)]
            ser.marker.size = 8
            ser.marker.format.fill.solid()
            ser.marker.format.fill.fore_color.rgb = RGBColor(*color)
            ser.marker.format.line.fill.background()
            ser.format.line.fill.background()
