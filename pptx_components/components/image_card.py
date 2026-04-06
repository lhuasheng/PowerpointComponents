from __future__ import annotations

from pathlib import Path
import warnings

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from pptx_components.base import Component, _resolve, add_rect, set_text_frame_margins
from pptx_components.theme import Theme


_VALID_MODES = ("contain", "stretch", "fit_width", "fit_height")


class ImageCard(Component):
    """Render an image card with optional caption and badge."""

    def __init__(
        self,
        image_path: str,
        caption: str | None = None,
        badge_text: str | None = None,
        mode: str = "contain",
        border_rgb: tuple[int, int, int] | None = None,
    ):
        if mode not in _VALID_MODES:
            raise ValueError(
                f"mode must be one of {list(_VALID_MODES)}; got {mode!r}"
            )
        self.image_path = image_path
        self.caption = caption
        self.badge_text = badge_text
        self.mode = mode
        self.border_rgb = border_rgb

    @property
    def min_height(self) -> float:
        return 1.45 if self.caption else 1.2

    def _remove_shape(self, shape) -> None:
        el = shape._element
        el.getparent().remove(el)

    def _apply_picture_border(self, pic) -> None:
        if self.border_rgb is None:
            pic.line.fill.background()
            return
        pic.line.fill.solid()
        pic.line.color.rgb = RGBColor(*self.border_rgb)
        pic.line.width = Pt(1)

    def _draw_missing_placeholder(self, slide, x: float, y: float, width: float, height: float,
                                  theme: Theme) -> None:
        ph = add_rect(slide, x, y, width, height, fill_rgb=theme.SURFACE_ALT, radius=0.03)
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"Image not found: {Path(self.image_path).name}"
        run.font.name = "Calibri"
        run.font.size = Pt(theme.CAPTION)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*theme.TEXT_MUTED)

    def _draw_badge(self, slide, x: float, y: float, width: float, theme: Theme) -> None:
        if not self.badge_text:
            return

        badge_h = 0.24
        est = 0.58 + (0.055 * len(self.badge_text))
        badge_w = min(max(1.0, est), max(1.0, width - 0.08))

        badge = add_rect(slide, x + 0.04, y + 0.04, badge_w, badge_h, fill_rgb=theme.ACCENT, radius=0.25)
        tf = badge.text_frame
        tf.clear()
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_text_frame_margins(tf, left=0.03, top=0.0, right=0.03, bottom=0.0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = self.badge_text
        run.font.name = "Calibri"
        font_size = theme.CAPTION
        available_w = max(badge_w - 0.06, 0.3)
        while font_size > 8:
            est_text_w = len(self.badge_text) * (font_size / 72.0) * 0.58
            if est_text_w <= available_w:
                break
            font_size -= 1
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)

    def render(
        self,
        slide,
        x: float,
        y: float,
        width: float,
        height: float,
        theme: Theme | None = None,
    ) -> None:
        t = _resolve(theme)
        pad = t.XS

        card = add_rect(slide, x, y, width, height, fill_rgb=t.SURFACE, radius=0.03)
        card.line.fill.solid()
        card.line.color.rgb = RGBColor(*(self.border_rgb if self.border_rgb is not None else t.SURFACE_ALT))
        card.line.width = Pt(1)

        caption_h = 0.24 if self.caption else 0.0
        caption_gap = t.XS if self.caption else 0.0
        image_x = x + pad
        image_y = y + pad
        image_w = max(width - 2 * pad, 0.1)
        image_h = max(height - (2 * pad) - caption_h - caption_gap, 0.1)

        path = Path(self.image_path)
        if not path.exists():
            warnings.warn(f"ImageCard missing image_path: {self.image_path}")
            self._draw_missing_placeholder(slide, image_x, image_y, image_w, image_h, t)
        elif self.mode == "stretch":
            pic = slide.shapes.add_picture(
                str(path), Inches(image_x), Inches(image_y), Inches(image_w), Inches(image_h)
            )
            self._apply_picture_border(pic)
        elif self.mode == "fit_width":
            pic = slide.shapes.add_picture(str(path), Inches(image_x), Inches(image_y), width=Inches(image_w))
            self._apply_picture_border(pic)
        elif self.mode == "fit_height":
            pic = slide.shapes.add_picture(str(path), Inches(image_x), Inches(image_y), height=Inches(image_h))
            self._apply_picture_border(pic)
        else:
            pic = slide.shapes.add_picture(str(path), Inches(image_x), Inches(image_y), width=Inches(image_w))
            if pic.height > Inches(image_h):
                self._remove_shape(pic)
                pic = slide.shapes.add_picture(
                    str(path), Inches(image_x), Inches(image_y), height=Inches(image_h)
                )

            target_left = Inches(image_x)
            target_top = Inches(image_y)
            target_width = Inches(image_w)
            target_height = Inches(image_h)
            pic.left = target_left + (target_width - pic.width) // 2
            pic.top = target_top + (target_height - pic.height) // 2
            self._apply_picture_border(pic)

        self._draw_badge(slide, image_x, image_y, image_w, t)

        if self.caption:
            cap = slide.shapes.add_textbox(
                Inches(x + pad),
                Inches(y + height - caption_h - pad),
                Inches(width - 2 * pad),
                Inches(caption_h),
            )
            tf = cap.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = self.caption
            run.font.name = "Calibri"
            run.font.bold = True
            run.font.size = Pt(t.CAPTION)
            run.font.color.rgb = RGBColor(*t.TEXT_SECONDARY)
