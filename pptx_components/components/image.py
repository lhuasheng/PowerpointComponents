from __future__ import annotations

from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from pptx_components.base import Component, _resolve
from pptx_components.theme import Theme


class ImageBlock(Component):
    """Render an image inside a bounding box.

    Args:
        image_path: Path to the source image.
        mode: "contain" (preserve ratio in box), "stretch" (fill box),
            "fit_width" (full width), or "fit_height" (full height).
        border_rgb: Optional border color tuple.
        border_width_pt: Border width when border_rgb is provided.
    """

    def __init__(
        self,
        image_path: str,
        mode: str = "contain",
        border_rgb: tuple[int, int, int] | None = None,
        border_width_pt: float = 1.0,
    ):
        if mode not in ("contain", "stretch", "fit_width", "fit_height"):
            raise ValueError(
                f"mode must be 'contain', 'stretch', 'fit_width', or 'fit_height'; got {mode!r}"
            )
        self.image_path = image_path
        self.mode = mode
        self.border_rgb = border_rgb
        self.border_width_pt = border_width_pt

    @property
    def min_height(self) -> float:
        return 1.8

    def _remove_shape(self, shape) -> None:
        el = shape._element
        el.getparent().remove(el)

    def _apply_border(self, shape) -> None:
        if self.border_rgb is None:
            shape.line.fill.background()
            return
        shape.line.fill.solid()
        shape.line.color.rgb = RGBColor(*self.border_rgb)
        shape.line.width = Pt(self.border_width_pt)

    def render(self, slide, x: float, y: float, width: float, height: float,
               theme: Theme | None = None) -> None:
        _resolve(theme)

        p = Path(self.image_path)
        if not p.exists():
            raise FileNotFoundError(f"image_path does not exist: {self.image_path}")

        if self.mode == "stretch":
            pic = slide.shapes.add_picture(
                str(p), Inches(x), Inches(y), Inches(width), Inches(height)
            )
            self._apply_border(pic)
            return

        if self.mode == "fit_width":
            pic = slide.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(width))
            self._apply_border(pic)
            return

        if self.mode == "fit_height":
            pic = slide.shapes.add_picture(str(p), Inches(x), Inches(y), height=Inches(height))
            self._apply_border(pic)
            return

        # contain: try fit-width first, then fallback to fit-height if it overflows.
        pic = slide.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(width))
        if pic.height > Inches(height):
            self._remove_shape(pic)
            pic = slide.shapes.add_picture(str(p), Inches(x), Inches(y), height=Inches(height))

        # Center image within target box.
        target_left = Inches(x)
        target_top = Inches(y)
        target_width = Inches(width)
        target_height = Inches(height)
        pic.left = target_left + (target_width - pic.width) // 2
        pic.top = target_top + (target_height - pic.height) // 2
        self._apply_border(pic)
