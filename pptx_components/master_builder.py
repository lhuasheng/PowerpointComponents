"""Master-driven slide builder — create decks from an existing PPTX template.

Use this when you want to keep a brand template's slide master (backgrounds,
logos, colour bars, fonts) while populating content programmatically.

Workflow::

    from pptx_components.master_builder import MasterPresentation
    import pptx_components as pc

    prs = MasterPresentation("brand_template.pptx")

    # Add a cover slide and fill its placeholders
    prs.add_slide("Cover Option - Generic (Default)", {
        0: "Q1 2026 Programme Review",
        11: "April 2026",
    })

    # Add a content slide and overlay a pptx_components chart
    slide = prs.add_slide("3_Title and Content", {0: "Revenue Trends"})
    slide.set_cursor(1.5).add(pc.BarChart(...), h=2.8)

    prs.save("output.pptx")
"""
from __future__ import annotations

import warnings
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from pptx_components.base import Component, _resolve
from pptx_components.layout import Row
from pptx_components.theme import Theme


# ── Single-slide editor ────────────────────────────────────────────────────

class MasterSlide:
    """Fluent editor for a single slide that inherits a master layout.

    Returned by :meth:`MasterPresentation.add_slide`.  Provides the same
    component-placement API as ``SlideBuilder`` (cursor, ``add``,
    ``add_row``, ``skip``) plus ``set_placeholder`` for populating the
    layout's text placeholders.
    """

    def __init__(
        self,
        slide,
        slide_w: float,
        slide_h: float,
        margin: float = 0.5,
        theme: Theme | None = None,
    ):
        self.slide = slide
        self.slide_w = slide_w
        self.slide_h = slide_h
        self.margin = margin
        self.theme = _resolve(theme)
        self.cursor_y: float = margin

    # ── Placeholder helpers ────────────────────────────────────────────────

    def set_placeholder(self, idx: int, text: str) -> "MasterSlide":
        """Set *text* in the placeholder at *idx*.  Returns self for chaining."""
        for ph in self.slide.placeholders:
            if ph.placeholder_format.idx == idx:
                tf = ph.text_frame
                # Clear runs within paragraphs to preserve paragraph formatting
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        r = run._r
                        r.getparent().remove(r)
                # Remove all paragraphs after the first
                while len(tf.paragraphs) > 1:
                    p = tf.paragraphs[1]._element
                    p.getparent().remove(p)
                # Set text on the first paragraph, preserving its formatting
                if tf.paragraphs:
                    tf.paragraphs[0].text = text
                else:
                    tf.text = text
                return self
        warnings.warn(
            f"Placeholder idx={idx} not found on this slide — skipping.",
            stacklevel=2,
        )
        return self

    def set_placeholders(self, mapping: dict[int, str]) -> "MasterSlide":
        """Set multiple placeholders at once.  *mapping* is ``{idx: text}``."""
        for idx, text in mapping.items():
            self.set_placeholder(idx, text)
        return self

    def hide_placeholder(self, idx: int) -> "MasterSlide":
        """Remove the placeholder shape at *idx* from the slide entirely."""
        sp_tree = self.slide.shapes._spTree
        for ph in list(self.slide.placeholders):
            if ph.placeholder_format.idx == idx:
                sp_tree.remove(ph._element)
                return self
        warnings.warn(
            f"Placeholder idx={idx} not found on this slide — skipping.",
            stacklevel=2,
        )
        return self

    # ── Component placement ────────────────────────────────────────────────

    def _content_width(self) -> float:
        return self.slide_w - 2 * self.margin

    def add(
        self,
        component: Component,
        x: float | None = None,
        y: float | None = None,
        w: float | None = None,
        h: float | None = None,
        theme: Theme | None = None,
    ) -> "MasterSlide":
        """Render *component* on the slide within the given bounding box (inches).

        Defaults mirror :class:`~pptx_components.SlideBuilder`:

        * ``x`` → ``margin``
        * ``w`` → ``slide_w - 2 * margin``
        * ``h`` → ``component.min_height_for(theme)``
        * ``y`` → current ``cursor_y``; cursor advances after render.

        Passing an explicit *y* pins the component without moving the cursor.
        """
        t = theme or self.theme
        resolved_x = x if x is not None else self.margin
        resolved_w = w if w is not None else self._content_width()
        resolved_h = h if h is not None else component.min_height_for(t)
        explicit_y = y is not None
        resolved_y = y if explicit_y else self.cursor_y

        component.render(self.slide, resolved_x, resolved_y, resolved_w, resolved_h, theme=t)

        if not explicit_y:
            gap = getattr(t, "SM", 0.2)
            self.cursor_y += resolved_h + gap

        return self

    def add_row(
        self,
        *components: Component,
        h: float | None = None,
        gap: float | None = None,
        weights: list[float] | None = None,
    ) -> "MasterSlide":
        """Wrap *components* in a :class:`~pptx_components.Row` and add at cursor."""
        row = Row(*components, gap=gap, weights=weights)
        return self.add(row, h=h)

    def skip(self, height: float) -> "MasterSlide":
        """Advance the cursor by *height* inches (manual spacing)."""
        self.cursor_y += height
        return self

    def set_cursor(self, y: float) -> "MasterSlide":
        """Manually position the vertical cursor to *y* inches."""
        self.cursor_y = y
        return self

    # ── Introspection ──────────────────────────────────────────────────────

    def list_placeholders(self) -> list[dict]:
        """Return info about every placeholder on this slide (useful for discovery)."""
        result = []
        for ph in self.slide.placeholders:
            fmt = ph.placeholder_format
            text = ph.text_frame.text if ph.has_text_frame else None
            result.append({
                "idx": fmt.idx,
                "type": str(fmt.type),
                "name": ph.name,
                "text": text,
            })
        return result


# ── Presentation builder ───────────────────────────────────────────────────

class MasterPresentation:
    """Build a slide deck using an existing PPTX file as the slide master.

    Opens *template_path* and, by default, removes all content slides so you
    start with an empty deck that retains the brand master (backgrounds,
    logos, colour bars, footer text, fonts).  New slides are added via
    :meth:`add_slide` by referencing layouts by name or index.

    Args:
        template_path: Path to the ``.pptx`` brand template.
        clear_slides:  Remove existing content slides (default ``True``).
        theme:         Default :class:`~pptx_components.Theme` used when
                       rendering pptx_components on slides.
        margin:        Default margin in inches for component placement.

    Example::

        prs = MasterPresentation("template.pptx")
        print(prs.layouts)               # see available layout names

        cover = prs.add_slide("Cover Option - Generic (Default)", {
            0: "My Report Title",
            11: "April 2026",
        })

        content = prs.add_slide("3_Title and Content", {0: "Results"})
        content.set_cursor(1.5).add(pc.BarChart(...))

        prs.save("my_report.pptx")
    """

    def __init__(
        self,
        template_path: str | Path,
        *,
        clear_slides: bool = True,
        theme: Theme | None = None,
        margin: float = 0.5,
    ):
        self._prs = Presentation(str(template_path))
        self._template_path = Path(template_path)
        self._theme = _resolve(theme)
        self._margin = margin
        self._slide_w = self._prs.slide_width.inches
        self._slide_h = self._prs.slide_height.inches

        # Build layout lookup: name → layout object, preserving order
        self._layouts: dict[str, object] = {}
        self._layout_list: list = []
        for master in self._prs.slide_masters:
            for layout in master.slide_layouts:
                self._layouts[layout.name] = layout
                self._layout_list.append(layout)

        if clear_slides:
            self._clear_existing_slides()

    # ── Internal ───────────────────────────────────────────────────────────

    def _clear_existing_slides(self) -> None:
        """Remove all existing slides while keeping the master and layouts."""
        prs_part = self._prs.part
        sldIdLst = self._prs.slides._sldIdLst
        # Collect rIds first, then drop each relationship (which also removes
        # the slide part from the package, avoiding duplicate-name warnings).
        rIds = [sldId.rId for sldId in list(sldIdLst)]
        for rId in rIds:
            prs_part.drop_rel(rId)
        for sldId in list(sldIdLst):
            sldIdLst.remove(sldId)

    # ── Public interface ───────────────────────────────────────────────────

    @property
    def layouts(self) -> list[str]:
        """Names of every available slide layout in the master."""
        return list(self._layouts.keys())

    @property
    def slide_width(self) -> float:
        """Template slide width in inches."""
        return self._slide_w

    @property
    def slide_height(self) -> float:
        """Template slide height in inches."""
        return self._slide_h

    def add_slide(
        self,
        layout: str | int,
        placeholders: dict[int, str] | None = None,
        theme: Theme | None = None,
    ) -> MasterSlide:
        """Add a new slide using the named (or 0-based indexed) layout.

        Args:
            layout:       Layout name (``str``) or 0-based index (``int``).
            placeholders: Optional ``{placeholder_idx: text}`` mapping to
                          pre-fill placeholder content immediately.
            theme:        Override the presentation-level theme for this slide.

        Returns:
            :class:`MasterSlide` — fluent single-slide editor.

        Raises:
            ValueError: If a string *layout* name is not found.
            IndexError: If an integer *layout* index is out of range.
        """
        if isinstance(layout, int):
            layout_obj = self._layout_list[layout]
        else:
            if layout not in self._layouts:
                available = "\n  ".join(self._layouts)
                raise ValueError(
                    f"Layout {layout!r} not found in template.\n"
                    f"Available layouts:\n  {available}"
                )
            layout_obj = self._layouts[layout]

        slide = self._prs.slides.add_slide(layout_obj)
        ms = MasterSlide(
            slide,
            slide_w=self._slide_w,
            slide_h=self._slide_h,
            margin=self._margin,
            theme=theme or self._theme,
        )
        if placeholders:
            ms.set_placeholders(placeholders)
        return ms

    def save(self, path: str | Path) -> str:
        """Save the presentation to *path*.

        Creates parent directories if needed.  Returns the absolute path
        as a string.
        """
        out = Path(path)
        out.parent.mkdir(parents=True, exist_ok=True)
        self._prs.save(str(out))
        return str(out.resolve())
