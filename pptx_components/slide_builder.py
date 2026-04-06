from __future__ import annotations

import warnings
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from pptx_components.base import Component, set_slide_background, _resolve
from pptx_components.delegation import GetAttr, get_first_attr
from pptx_components.layout import Row
from pptx_components.theme import Theme, get_theme


@dataclass(frozen=True)
class LayoutIssue:
    """Structured representation of an overflow issue found during layout."""

    slide_number: int
    severity: str
    component_name: str
    message: str
    y: float
    h: float
    safe_bottom: float
    overflow: float


class SlideBuilder(GetAttr):
    """Clean API for composing a single slide from components.

    Tracks a vertical cursor so callers can `add()` components sequentially
    without manually computing y-offsets.
    """

    _default = "theme"

    def __init__(
        self,
        prs: Presentation,
        theme: Theme | None = None,
        validate: bool = False,
        strict: bool = False,
    ):
        self._prs = prs
        self.theme = _resolve(theme)
        # strict mode always enables validation.
        self.validate = validate or strict
        self.strict = strict
        self._slide_h = prs.slide_height.inches
        # Use blank slide layout (index 6)
        blank_layout = prs.slide_layouts[6]
        self.slide = prs.slides.add_slide(blank_layout)
        self.slide_number = len(prs.slides)
        self.layout_issues: list[LayoutIssue] = []
        bg_image = get_first_attr(self.theme, "BG_IMAGE")
        if bg_image:
            bg_path = Path(bg_image)
            if bg_path.exists():
                self.slide.shapes.add_picture(
                    str(bg_path),
                    Inches(0),
                    Inches(0),
                    width=prs.slide_width,
                    height=prs.slide_height,
                )
            else:
                warnings.warn(f"Background image not found, using BG color: {bg_image}")
                set_slide_background(self.slide, self.BG)
        else:
            set_slide_background(self.slide, self.BG)
        self.cursor_y: float = self.MARGIN

        logo_path = get_first_attr(self.theme, "LOGO_PATH", "logo_path")
        if logo_path:
            logo_x = get_first_attr(self.theme, "LOGO_X", "logo_x")
            logo_y = get_first_attr(self.theme, "LOGO_Y", "logo_y")
            logo_w = get_first_attr(self.theme, "LOGO_W", "logo_w")
            if logo_x is not None and logo_y is not None and logo_w is not None:
                self.set_logo(str(logo_path), float(logo_x), float(logo_y), float(logo_w))

    # ── Internal ───────────────────────────────────────────────────────────

    def _content_width(self) -> float:
        return self.SLIDE_W - 2 * self.MARGIN

    def _safe_bottom(self) -> float:
        return self._slide_h - self.MARGIN

    def _handle_overflow(
        self,
        *,
        component: Component,
        y: float,
        h: float,
        allow_overflow: bool,
    ) -> None:
        if not self.validate or allow_overflow:
            return

        safe_bottom = self._safe_bottom()
        bottom = y + h
        if bottom <= safe_bottom:
            return

        overflow = bottom - safe_bottom
        comp_name = component.__class__.__name__
        message = (
            f"SlideBuilder overflow for {comp_name}: y+h={bottom:.2f}\" exceeds "
            f"safe bottom {safe_bottom:.2f}\" by {overflow:.2f}\". "
            f"Placement y={y:.2f}\", h={h:.2f}\". "
            "Use allow_overflow=True to bypass for intentional overlays."
        )

        severity = "error" if self.strict else "warning"
        self.layout_issues.append(
            LayoutIssue(
                slide_number=self.slide_number,
                severity=severity,
                component_name=comp_name,
                message=message,
                y=y,
                h=h,
                safe_bottom=safe_bottom,
                overflow=overflow,
            )
        )

        if self.strict:
            raise ValueError(message)

        warnings.warn(message, stacklevel=2)

    # ── Public API ─────────────────────────────────────────────────────────

    def add(self, component: Component,
            x: float | None = None,
            y: float | None = None,
            w: float | None = None,
            h: float | None = None,
            allow_overflow: bool = False) -> "SlideBuilder":
        """Render a component on the slide.

        Defaults:
          x = theme.MARGIN
          w = SLIDE_W - 2*MARGIN
          h = component.min_height
          y = cursor_y (auto-advances after render)

        Passing an explicit y overrides the cursor for this call only
        (cursor is NOT advanced when y is explicitly provided).

        When validation is enabled, placement is checked against safe bottom
        bounds (slide height - margin). Use allow_overflow=True to bypass.
        """
        t = self.theme
        resolved_x = x if x is not None else t.MARGIN
        resolved_w = w if w is not None else self._content_width()
        resolved_h = h if h is not None else component.min_height_for(t)

        explicit_y = y is not None
        resolved_y = y if explicit_y else self.cursor_y

        self._handle_overflow(
            component=component,
            y=resolved_y,
            h=resolved_h,
            allow_overflow=allow_overflow,
        )

        component.render(self.slide, resolved_x, resolved_y, resolved_w, resolved_h, theme=t)

        if not explicit_y:
            self.cursor_y += resolved_h + t.SM

        return self  # fluent API

    def add_full(self, component: Component,
                 h: float | None = None,
                 allow_overflow: bool = False) -> "SlideBuilder":
        """Add a component spanning the full content width at the current cursor."""
        return self.add(component, h=h, allow_overflow=allow_overflow)

    def add_row(self, *components: Component,
                h: float | None = None,
                gap: float | None = None,
                weights: list[float] | None = None,
                allow_overflow: bool = False) -> "SlideBuilder":
        """Wrap components in a Row and add at the current cursor."""
        row = Row(*components, gap=gap, weights=weights)
        return self.add(row, h=h, allow_overflow=allow_overflow)

    def skip(self, height: float) -> "SlideBuilder":
        """Advance the cursor by a fixed amount (manual spacing)."""
        self.cursor_y += height
        return self

    def set_cursor(self, y: float) -> "SlideBuilder":
        """Manually position the cursor."""
        self.cursor_y = y
        return self

    def set_logo(self, path: str, x: float, y: float, w: float) -> "SlideBuilder":
        """Place a logo image on this slide at inch-based coordinates."""
        p = Path(path)
        if not p.exists():
            warnings.warn(f"Logo not found, skipping: {path}")
            return self

        self.slide.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w))
        return self
