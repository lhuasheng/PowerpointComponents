from __future__ import annotations
from abc import ABC, abstractmethod

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from pptx_components.theme import Theme, get_theme


# ── Abstract base ──────────────────────────────────────────────────────────

class Component(ABC):
    @abstractmethod
    def render(self, slide, x: float, y: float, width: float, height: float,
               theme: Theme | None = None) -> None:
        """Render the component onto *slide* within the given bounding box (inches)."""

    @property
    @abstractmethod
    def min_height(self) -> float:
        """Minimum height in inches this component needs to render correctly."""

    def min_height_for(self, theme: Theme | None = None) -> float:
        """Theme-aware min-height hook for cascading layout measurements.

        Components with truly theme-dependent sizing can override this method.
        """
        _resolve(theme)
        return self.min_height


# ── Module-level helpers (stateless, no self) ──────────────────────────────

def _resolve(theme: Theme | None) -> Theme:
    return theme if theme is not None else get_theme()


def resolve_theme(*themes: Theme | None) -> Theme:
    """Resolve a theme by precedence (first non-None wins), else global default."""
    for theme in themes:
        if theme is not None:
            return theme
    return get_theme()


def add_rect(slide, x: float, y: float, w: float, h: float,
             fill_rgb: tuple[int, int, int] | None = None,
             radius: float = 0.0):
    """Add a rectangle (or rounded rectangle) to the slide. Returns the shape."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches

    if radius > 0:
        # MSO shape type 5 = rounded rectangle
        shape = slide.shapes.add_shape(
            5,  # MSO_SHAPE_TYPE.ROUNDED_RECTANGLE
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        # Set corner radius via adjustment value (0–50000 EMU scale; 50000 = fully round)
        # A radius of 0.05 inch relative to min dimension gives a subtle rounding
        min_dim_emu = min(Inches(w), Inches(h))
        adj = int(radius * 914400 / min_dim_emu * 100000)
        adj = max(0, min(50000, adj))
        shape.adjustments[0] = adj / 100000
    else:
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(x), Inches(y), Inches(w), Inches(h)
        )

    apply_no_line(shape)
    if fill_rgb is not None:
        apply_fill(shape, fill_rgb)
    else:
        apply_no_fill(shape)

    return shape


def apply_fill(shape, rgb: tuple[int, int, int]) -> None:
    """Apply a solid fill color to a shape."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb)


def apply_no_fill(shape) -> None:
    """Make a shape transparent (no fill)."""
    shape.fill.background()


def apply_no_line(shape) -> None:
    """Remove the border line from a shape."""
    shape.line.fill.background()


def set_font(run, size: int, bold: bool = False, italic: bool = False,
             color_rgb: tuple[int, int, int] | None = None,
             font_name: str = "Calibri") -> None:
    """Style a text run."""
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color_rgb is not None:
        run.font.color.rgb = RGBColor(*color_rgb)


def set_para_align(para, alignment) -> None:
    """Set paragraph alignment (PP_ALIGN.LEFT / CENTER / RIGHT)."""
    para.alignment = alignment


def set_text_frame_margins(tf, left: float = 0.05, top: float = 0.05,
                            right: float = 0.05, bottom: float = 0.05) -> None:
    """Set internal text frame margins in inches."""
    tf.margin_left = Inches(left)
    tf.margin_top = Inches(top)
    tf.margin_right = Inches(right)
    tf.margin_bottom = Inches(bottom)


def add_accent_bar(slide, x: float, y: float, h: float,
                   theme: Theme, width: float = 0.05) -> None:
    """Draw a thin vertical accent-colored bar on the left edge of a component."""
    bar = add_rect(slide, x, y, width, h, fill_rgb=theme.ACCENT)
    return bar


def add_text_box(slide, x: float, y: float, w: float, h: float,
                 text: str, size: int, bold: bool = False, italic: bool = False,
                 color_rgb: tuple[int, int, int] | None = None,
                 alignment=PP_ALIGN.LEFT,
                 font_name: str = "Calibri",
                 word_wrap: bool = True) -> None:
    """Add a standalone text box with a single run of styled text."""
    from pptx.util import Inches
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    set_text_frame_margins(tf, 0, 0, 0, 0)

    para = tf.paragraphs[0]
    para.alignment = alignment
    run = para.add_run()
    run.text = text
    set_font(run, size, bold=bold, italic=italic,
             color_rgb=color_rgb, font_name=font_name)
    return txBox


def set_slide_background(slide, rgb: tuple[int, int, int]) -> None:
    """Fill the slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)
