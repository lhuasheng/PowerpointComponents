from __future__ import annotations

import argparse
import os
import re
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


EMU_PER_INCH = 914400.0


@dataclass
class ReverseWarning:
    slide_index: int
    message: str


@dataclass
class ReverseResult:
    script_path: Path
    assets_dir: Path
    warnings: list[ReverseWarning]


def _to_inches(emu_value: int) -> float:
    return round(float(emu_value) / EMU_PER_INCH, 3)


def _py_str(value: str) -> str:
    return repr(value)


def _safe_name(value: str) -> str:
    return re.sub(r"[^a-zA-Z0-9_-]+", "_", value).strip("_") or "item"


def _shape_sort_key(shape) -> tuple[int, int, int, int]:
    return (int(shape.top), int(shape.left), int(shape.height), int(shape.width))


class PresentationReverser:
    """Convert an edited PPTX into a runnable SlideBuilder script.

    Uses best-effort mapping for common PowerPoint objects and falls back to
    positioned placeholder shapes for unsupported objects.
    """

    _PIE_MARKERS = ("PIE", "DOUGHNUT")
    _LINE_MARKERS = ("LINE",)
    _BAR_MARKERS = ("BAR", "COLUMN")

    @staticmethod
    def _chart_type_name(chart) -> str:
        return str(getattr(chart, "chart_type", "")).upper()

    def _chart_to_component_call(self, shape, x: float, y: float, w: float, h: float) -> str | None:
        if not getattr(shape, "has_chart", False):
            return None
        chart = shape.chart
        chart_type_name = self._chart_type_name(chart)

        title: str | None = None
        try:
            if chart.has_title and chart.chart_title and chart.chart_title.text_frame:
                raw_title = chart.chart_title.text_frame.text.strip()
                title = raw_title if raw_title else None
        except Exception:
            title = None

        categories: list[str] = []
        try:
            categories = [str(c.label) for c in chart.plots[0].categories]
        except Exception:
            categories = []

        series_dict: dict[str, list[float]] = {}
        try:
            for i, series in enumerate(chart.series, start=1):
                name = (series.name or "").strip() or f"Series {i}"
                values = [float(v) if v is not None else 0.0 for v in series.values]
                series_dict[name] = values
        except Exception:
            series_dict = {}

        if not categories and series_dict:
            max_len = max(len(v) for v in series_dict.values())
            categories = [f"Item {i}" for i in range(1, max_len + 1)]

        if not categories or not series_dict:
            return None

        if any(marker in chart_type_name for marker in self._PIE_MARKERS):
            first_series_name = next(iter(series_dict.keys()))
            values = series_dict[first_series_name]
            return (
                f"    b.add(pc.PieChart(categories={repr(categories)}, values={repr(values)}, title={repr(title)}), "
                f"x={x}, y={y}, w={w}, h={h})"
            )

        if any(marker in chart_type_name for marker in self._LINE_MARKERS):
            return (
                f"    b.add(pc.LineChart(categories={repr(categories)}, series={repr(series_dict)}, title={repr(title)}), "
                f"x={x}, y={y}, w={w}, h={h})"
            )

        if any(marker in chart_type_name for marker in self._BAR_MARKERS):
            return (
                f"    b.add(pc.BarChart(categories={repr(categories)}, series={repr(series_dict)}, title={repr(title)}), "
                f"x={x}, y={y}, w={w}, h={h})"
            )

        return None

    @staticmethod
    def _shape_label(shape) -> str:
        shape_type = str(getattr(shape, "shape_type", "UNKNOWN"))
        shape_name = (getattr(shape, "name", "") or "").strip()
        if shape_name:
            return f"{shape_type}: {shape_name}"
        return shape_type

    def __init__(self, pptx_path: str | Path):
        self.pptx_path = Path(pptx_path).expanduser().resolve()
        if not self.pptx_path.exists():
            raise FileNotFoundError(f"PPTX not found: {self.pptx_path}")
        self.prs = Presentation(str(self.pptx_path))

    def reverse_to_script(
        self,
        output_script_path: str | Path,
        assets_dir: str | Path | None = None,
        strict: bool = False,
    ) -> ReverseResult:
        output_script = Path(output_script_path).expanduser().resolve()
        output_script.parent.mkdir(parents=True, exist_ok=True)

        if assets_dir is None:
            assets = output_script.parent / f"{output_script.stem}_assets"
        else:
            assets = Path(assets_dir).expanduser().resolve()
        assets.mkdir(parents=True, exist_ok=True)

        warnings: list[ReverseWarning] = []
        picture_assets_dir = assets / "images"
        picture_assets_dir.mkdir(parents=True, exist_ok=True)

        component_lines_by_slide: list[list[str]] = []

        for slide_idx, slide in enumerate(self.prs.slides, start=1):
            lines: list[str] = []
            unsupported_shape_count = 0
            picture_counter = 0

            shapes = sorted(slide.shapes, key=_shape_sort_key)
            for shape in shapes:
                x = _to_inches(shape.left)
                y = _to_inches(shape.top)
                w = _to_inches(shape.width)
                h = _to_inches(shape.height)

                # Picture -> ImageBlock with extracted binary payload.
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    ext = image.ext or "png"
                    picture_counter += 1
                    img_name = f"slide_{slide_idx:03d}_img_{picture_counter:03d}.{_safe_name(ext)}"
                    img_path = picture_assets_dir / img_name
                    img_path.write_bytes(image.blob)
                    lines.append(
                        f"    b.add(pc.ImageBlock(str(ASSETS_DIR / 'images' / {_py_str(img_name)})), x={x}, y={y}, w={w}, h={h})"
                    )
                    continue

                chart_line = self._chart_to_component_call(shape, x, y, w, h)
                if chart_line:
                    lines.append(chart_line)
                    continue

                # Table -> DataTable
                if getattr(shape, "has_table", False):
                    table = shape.table
                    if table.rows and table.columns:
                        headers = [table.cell(0, c).text.strip() for c in range(len(table.columns))]
                        rows = [
                            [table.cell(r, c).text.strip() for c in range(len(table.columns))]
                            for r in range(1, len(table.rows))
                        ]
                        lines.append(
                            f"    b.add(pc.DataTable(headers={repr(headers)}, rows={repr(rows)}), x={x}, y={y}, w={w}, h={h})"
                        )
                        continue

                # Text frame -> TitleBlock/ListBlock/TextCard
                if getattr(shape, "has_text_frame", False):
                    paragraphs = [p for p in shape.text_frame.paragraphs if p.text and p.text.strip()]
                    text_lines = [p.text.strip() for p in paragraphs]
                    if text_lines:
                        # Heuristic title: prominent top-of-slide text with short line count.
                        if y <= 1.8 and len(text_lines) <= 3 and len(" ".join(text_lines)) <= 180:
                            title = text_lines[0]
                            subtitle = " ".join(text_lines[1:]) if len(text_lines) > 1 else None
                            lines.append(
                                f"    b.add(pc.TitleBlock({_py_str(title)}, subtitle={repr(subtitle)}), x={x}, y={y}, w={w}, h={h})"
                            )
                            continue

                        list_like = any(p.level > 0 for p in paragraphs) or len(text_lines) >= 3
                        if list_like:
                            clean_items = [re.sub(r"^[\u2022\-\*\d\.\)\s]+", "", t).strip() for t in text_lines]
                            clean_items = [t for t in clean_items if t]
                            if clean_items:
                                lines.append(
                                    f"    b.add(pc.ListBlock({repr(clean_items)}), x={x}, y={y}, w={w}, h={h})"
                                )
                                continue

                        body_text = "\n".join(text_lines)
                        lines.append(
                            f"    b.add(pc.TextCard(body={_py_str(body_text)}), x={x}, y={y}, w={w}, h={h})"
                        )
                        continue

                    # Empty placeholders are ignored to avoid false fallback triggers.
                    continue

                unsupported_shape_count += 1
                label = self._shape_label(shape)
                text_value = ""
                if getattr(shape, "has_text_frame", False):
                    try:
                        text_value = shape.text_frame.text.strip()
                    except Exception:
                        text_value = ""

                lines.append(
                    f"    _add_unknown_shape_placeholder(b.slide, x={x}, y={y}, w={w}, h={h}, label={_py_str(label)}, text={_py_str(text_value)})"
                )

            if unsupported_shape_count:
                warnings.append(
                    ReverseWarning(
                        slide_index=slide_idx,
                        message=(
                            f"Slide {slide_idx}: rendered {unsupported_shape_count} unsupported shape(s) "
                            "as positioned placeholders."
                        ),
                    )
                )
                if strict:
                    raise ValueError(warnings[-1].message)

            component_lines_by_slide.append(lines)

        script = self._generate_script(
            output_script=output_script,
            assets_dir=assets,
            component_lines_by_slide=component_lines_by_slide,
            warnings=warnings,
        )
        output_script.write_text(script, encoding="utf-8")

        return ReverseResult(script_path=output_script, assets_dir=assets, warnings=warnings)

    def _generate_script(
        self,
        output_script: Path,
        assets_dir: Path,
        component_lines_by_slide: list[list[str]],
        warnings: list[ReverseWarning],
    ) -> str:
        rel_assets = Path(os.path.relpath(assets_dir, output_script.parent))
        slide_w = _to_inches(self.prs.slide_width)
        slide_h = _to_inches(self.prs.slide_height)

        warning_lines = [
            f"# - Slide {w.slide_index}: {w.message}" for w in warnings
        ] or ["# - none"]

        lines: list[str] = [
            "from __future__ import annotations",
            "",
            "from pathlib import Path",
            "",
            "from pptx.dml.color import RGBColor",
            "from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE",
            "from pptx.enum.text import PP_ALIGN",
            "from pptx import Presentation",
            "from pptx.util import Inches, Pt",
            "import pptx_components as pc",
            "",
            f"ASSETS_DIR = Path(__file__).resolve().parent / {repr(rel_assets.as_posix())}",
            f"SLIDE_W = {slide_w}",
            f"SLIDE_H = {slide_h}",
            "",
            "",
            "def _add_unknown_shape_placeholder(slide, x: float, y: float, w: float, h: float, label: str, text: str | None = None) -> None:",
            "    shape = slide.shapes.add_shape(",
            "        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,",
            "        Inches(x),",
            "        Inches(y),",
            "        Inches(max(w, 0.2)),",
            "        Inches(max(h, 0.2)),",
            "    )",
            "    shape.fill.solid()",
            "    shape.fill.fore_color.rgb = RGBColor(245, 245, 245)",
            "    shape.line.color.rgb = RGBColor(140, 140, 140)",
            "    shape.line.width = Pt(1)",
            "    tf = shape.text_frame",
            "    tf.clear()",
            "    p = tf.paragraphs[0]",
            "    p.alignment = PP_ALIGN.LEFT",
            "    run = p.add_run()",
            "    run.text = f'[Unsupported] {label}'",
            "    run.font.size = Pt(9)",
            "    run.font.bold = True",
            "    run.font.color.rgb = RGBColor(70, 70, 70)",
            "    if text:",
            "        p2 = tf.add_paragraph()",
            "        p2.alignment = PP_ALIGN.LEFT",
            "        run2 = p2.add_run()",
            "        run2.text = text[:240]",
            "        run2.font.size = Pt(8)",
            "        run2.font.color.rgb = RGBColor(90, 90, 90)",
            "",
            "# Auto-generated by pptx_components.reverse.",
            f"# Source PPTX: {self.pptx_path.name}",
            "# Mapping warnings:",
            *warning_lines,
            "",
        ]

        for idx, comp_lines in enumerate(component_lines_by_slide, start=1):
            lines.append(f"def add_slide_{idx:03d}(prs: Presentation) -> None:")
            lines.append("    b = pc.SlideBuilder(prs, theme=pc.LightTheme())")
            if comp_lines:
                lines.extend(comp_lines)
            else:
                lines.append("    # Slide contained no detectable mappable content.")
            lines.append("")

        lines.extend(
            [
                "def build(output_path: str | Path = 'reverse_generated.pptx') -> Path:",
                "    prs = Presentation()",
                "    prs.slide_width = Inches(SLIDE_W)",
                "    prs.slide_height = Inches(SLIDE_H)",
                "",
            ]
        )

        for idx in range(1, len(component_lines_by_slide) + 1):
            lines.append(f"    add_slide_{idx:03d}(prs)")

        lines.extend(
            [
                "    out = Path(output_path).resolve()",
                "    out.parent.mkdir(parents=True, exist_ok=True)",
                "    prs.save(str(out))",
                "    return out",
                "",
                "",
                "if __name__ == '__main__':",
                "    out = build()",
                "    print(f'Generated PPTX: {out}')",
            ]
        )

        return "\n".join(lines) + "\n"


def reverse_pptx_to_script(
    pptx_path: str | Path,
    output_script_path: str | Path,
    assets_dir: str | Path | None = None,
    strict: bool = False,
) -> ReverseResult:
    reverser = PresentationReverser(pptx_path)
    return reverser.reverse_to_script(
        output_script_path=output_script_path,
        assets_dir=assets_dir,
        strict=strict,
    )


def _build_arg_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Reverse an edited PPTX into a SlideBuilder-style Python script.",
    )
    parser.add_argument("pptx", help="Path to source .pptx")
    parser.add_argument(
        "-o",
        "--output-script",
        default="reverse_generated.py",
        help="Output Python script path (default: reverse_generated.py)",
    )
    parser.add_argument(
        "--assets-dir",
        default=None,
        help="Assets directory for extracted images",
    )
    parser.add_argument(
        "--strict",
        action="store_true",
        help="Fail instead of fallback when unsupported shapes are found",
    )
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = _build_arg_parser()
    args = parser.parse_args(argv)
    result = reverse_pptx_to_script(
        pptx_path=args.pptx,
        output_script_path=args.output_script,
        assets_dir=args.assets_dir,
        strict=args.strict,
    )
    print(f"Generated script: {result.script_path}")
    print(f"Assets directory: {result.assets_dir}")
    if result.warnings:
        print("Warnings:")
        for warning in result.warnings:
            print(f"  - Slide {warning.slide_index}: {warning.message}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())